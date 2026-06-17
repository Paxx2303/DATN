from pptx import Presentation
from collections import Counter
import os

prs = Presentation(r'C:\Using\DATN\BaoVe_NguyenQuocNam_221220938.pptx')
sizes = Counter()
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        sizes[run.font.size] += 1
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                sizes[run.font.size] += 1

print('Font distribution:')
for sz, cnt in sorted(sizes.items()):
    print(f'  Pt({int(sz/12700)}) x{cnt}')

mb = os.path.getsize(r'C:\Using\DATN\BaoVe_NguyenQuocNam_221220938.pptx') / 1024 / 1024
print(f'File: {mb:.2f} MB, {len(prs.slides)} slides')
