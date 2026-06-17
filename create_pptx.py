# -*- coding: utf-8 -*-
"""Create BaoVe_NguyenQuocNam_221220938.pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── constants ──────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# colors
NAVY   = RGBColor(0x0D, 0x1F, 0x4E)   # dark navy (primary)
BLUE   = RGBColor(0x44, 0x72, 0xC4)   # accent blue
ORANGE = RGBColor(0xED, 0x7D, 0x31)   # accent orange
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)   # light gray bg
DGRAY  = RGBColor(0x44, 0x54, 0x6A)   # dark gray text
MGRAY  = RGBColor(0xBF, 0xBF, 0xBF)   # mid gray
TBLHDR = RGBColor(0x0D, 0x1F, 0x4E)   # table header (navy)
TBLALT = RGBColor(0xE8, 0xEF, 0xF8)   # table alt row
GREEN  = RGBColor(0x38, 0x8E, 0x3C)   # green for pass

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # blank

# ── helpers ────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             font_size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font_name='Calibri', italic=False,
             wrap=True, margin=Inches(0.1)):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    tf.margin_left   = margin
    tf.margin_right  = margin
    tf.margin_top    = margin
    tf.margin_bottom = margin
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    run.font.color.rgb = color
    return txb

def add_multiline(slide, lines, x, y, w, h,
                  font_size=Pt(13), bold=False, color=DGRAY,
                  align=PP_ALIGN.LEFT, font_name='Calibri',
                  line_spacing=None):
    """lines: list of (text, bold, color, size, italic) OR just strings"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left   = Inches(0.05)
    tf.margin_right  = Inches(0.05)
    tf.margin_top    = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, b, c, sz, it = item, bold, color, font_size, False
        else:
            txt = item[0]
            b   = item[1] if len(item)>1 else bold
            c   = item[2] if len(item)>2 else color
            sz  = item[3] if len(item)>3 else font_size
            it  = item[4] if len(item)>4 else False
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = txt
        run.font.size = sz
        run.font.bold = b
        run.font.italic = it
        run.font.name = font_name
        run.font.color.rgb = c
    return txb

def nav_bar(slide, title_text, slide_num, total=24):
    """Top navy header bar with title"""
    add_rect(slide, 0, 0, W, Inches(0.75), fill=NAVY)
    add_text(slide, title_text,
             Inches(0.3), Inches(0.1), Inches(11.5), Inches(0.6),
             font_size=Pt(20), bold=True, color=WHITE,
             font_name='Calibri Light')
    # slide number top right
    add_text(slide, f"{slide_num}/{total}",
             Inches(12.3), Inches(0.1), Inches(0.8), Inches(0.5),
             font_size=Pt(11), color=RGBColor(0xCC, 0xCC, 0xFF),
             align=PP_ALIGN.RIGHT)

def footer_bar(slide):
    """Bottom orange accent strip"""
    add_rect(slide, 0, H - Inches(0.25), W, Inches(0.25), fill=ORANGE)

def content_slide(slide):
    """White background content slide base"""
    add_rect(slide, 0, 0, W, H, fill=WHITE)
    footer_bar(slide)

def section_slide(slide, section_name):
    """Full dark slide as section divider"""
    add_rect(slide, 0, 0, W, H, fill=NAVY)
    add_rect(slide, 0, H - Inches(0.25), W, Inches(0.25), fill=ORANGE)
    add_text(slide, section_name,
             Inches(1), Inches(2.5), Inches(11), Inches(2),
             font_size=Pt(40), bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, font_name='Calibri Light')

def add_table(slide, headers, rows,
              x=Inches(0.4), y=Inches(1.0), w=Inches(12.5), h=None,
              col_widths=None, font_size=Pt(11), hdr_size=Pt(11)):
    """Add a styled table"""
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    ncols = len(headers)
    nrows = len(rows) + 1
    if h is None:
        h = Inches(0.35) * nrows + Inches(0.1)

    tbl = slide.shapes.add_table(nrows, ncols, x, y, w, h).table
    tbl.first_row = True

    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    def set_cell(cell, text, bold=False, bg=None, fg=WHITE, sz=Pt(11), align=PP_ALIGN.LEFT, italic=False):
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left   = Inches(0.05)
        tf.margin_right  = Inches(0.05)
        tf.margin_top    = Inches(0.03)
        tf.margin_bottom = Inches(0.03)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = sz
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = 'Calibri'
        run.font.color.rgb = fg
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    # header row
    for j, hdr in enumerate(headers):
        set_cell(tbl.cell(0, j), hdr, bold=True, bg=TBLHDR, fg=WHITE, sz=hdr_size, align=PP_ALIGN.CENTER)

    # data rows
    for i, row in enumerate(rows):
        bg = TBLALT if i % 2 == 0 else WHITE
        for j, val in enumerate(row):
            b = False
            fg = DGRAY
            it = False
            if isinstance(val, tuple):
                txt, b, fg = val[0], val[1] if len(val)>1 else False, val[2] if len(val)>2 else DGRAY
            else:
                txt = str(val)
            if txt.startswith('✅'):
                fg = GREEN
                b = True
            set_cell(tbl.cell(i+1, j), txt, bold=b, bg=bg, fg=fg, sz=font_size)

    return tbl

# ══════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)

# full navy bg
add_rect(sl, 0, 0, W, H, fill=NAVY)
# orange bottom bar
add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), fill=ORANGE)
# blue left accent
add_rect(sl, 0, 0, Inches(0.12), H, fill=ORANGE)

# university top
add_text(sl, "ĐẠI HỌC GIAO THÔNG VẬN TẢI",
         Inches(0.3), Inches(0.3), Inches(12.5), Inches(0.55),
         font_size=Pt(16), bold=False, color=RGBColor(0xCC,0xD5,0xFF),
         align=PP_ALIGN.CENTER, font_name='Calibri Light')

# white card in center
add_rect(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.8),
         fill=RGBColor(0xFF,0xFF,0xFF))

# main title inside card
add_text(sl, "XÂY DỰNG HỆ THỐNG NHẬN DIỆN VẬT THỂ\nQUA CAMERA MẮT CÁ",
         Inches(0.7), Inches(1.5), Inches(11.9), Inches(1.8),
         font_size=Pt(30), bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, font_name='Calibri Light')

# subtitle
add_text(sl, "ĐỒ ÁN TỐT NGHIỆP — CNTT1-K63",
         Inches(0.7), Inches(3.3), Inches(11.9), Inches(0.5),
         font_size=Pt(15), bold=False, color=ORANGE,
         align=PP_ALIGN.CENTER, font_name='Calibri')

# divider line inside card
add_rect(sl, Inches(2), Inches(3.8), Inches(9.3), Inches(0.02), fill=MGRAY)

# info below divider
info_lines = [
    ("Sinh viên:  Nguyễn Quốc Nam  ·  MSV 221220938  ·  Lớp CNTT1-K63", False, NAVY, Pt(13), False),
    ("GVHD:  TS. Nguyễn Đức Dư", False, NAVY, Pt(13), False),
    ("Bộ môn Công nghệ Thông tin  ·  Hà Nội, 2026", False, DGRAY, Pt(12), True),
]
add_multiline(sl, info_lines, Inches(0.7), Inches(3.9), Inches(11.9), Inches(1.0),
              align=PP_ALIGN.CENTER)

# bottom decoration
add_text(sl, "CNTT1-K63  ·  Khóa 63  ·  Hệ chính quy  ·  Ngành Công nghệ Thông tin",
         Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.4),
         font_size=Pt(10), color=RGBColor(0xAA,0xAA,0xFF),
         align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 2 — Đặt Vấn Đề
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Đặt Vấn Đề", 2)

# left column - thực trạng
add_rect(sl, Inches(0.3), Inches(0.9), Inches(6.1), Inches(2.3), fill=RGBColor(0xE8,0xEF,0xF8))
add_text(sl, "Thực Trạng Giao Thông",
         Inches(0.4), Inches(0.95), Inches(5.9), Inches(0.45),
         font_size=Pt(13), bold=True, color=NAVY)
add_multiline(sl, [
    ("• 7,8 triệu ô tô & 73 triệu xe máy (Tổng cục Thống kê, 2024)", True, NAVY, Pt(12)),
    ("• Năm 2023: hơn 10.000 vụ tai nạn giao thông nghiêm trọng", False, DGRAY, Pt(12)),
    ("• Camera CCTV truyền thống: giám sát thủ công, không phân tích tự động", False, DGRAY, Pt(12)),
], Inches(0.4), Inches(1.45), Inches(5.9), Inches(1.6))

# right column - fisheye advantage
add_rect(sl, Inches(6.7), Inches(0.9), Inches(6.2), Inches(2.3), fill=RGBColor(0xFF, 0xF3, 0xE0))
add_text(sl, "Tại Sao Camera Fisheye?",
         Inches(6.8), Inches(0.95), Inches(6.0), Inches(0.45),
         font_size=Pt(13), bold=True, color=ORANGE)
add_multiline(sl, [
    ("• Góc nhìn 180°–220° — một camera bao phủ toàn bộ ngã tư", True, RGBColor(0x8B,0x45,0x13), Pt(12)),
    ("• Tiết kiệm chi phí lắp đặt so với 4 camera thường", False, DGRAY, Pt(12)),
    ("• Nhưng gây barrel distortion — AI thông thường không áp dụng trực tiếp", False, DGRAY, Pt(12)),
], Inches(6.8), Inches(1.45), Inches(6.0), Inches(1.6))

# bottom - challenges
add_rect(sl, Inches(0.3), Inches(3.35), Inches(12.6), Inches(0.45), fill=NAVY)
add_text(sl, "Thách Thức Kỹ Thuật Đặc Thù",
         Inches(0.5), Inches(3.37), Inches(12.0), Inches(0.4),
         font_size=Pt(14), bold=True, color=WHITE)

challenges = [
    ("Bbox\nkém phù hợp", "Đối tượng ở vùng biên ảnh méo không vừa với bbox chữ nhật"),
    ("Đối tượng\nnhỏ", "Người đi bộ xa tâm ảnh chỉ 10–30 pixel chiều cao"),
    ("Kích thước\nphi tuyến", "Biến đổi theo vị trí trong ảnh fisheye"),
    ("Dữ liệu\nkhan hiếm", "Ảnh fisheye có nhãn chất lượng cao còn ít"),
]
for i, (title, desc) in enumerate(challenges):
    x = Inches(0.3 + i * 3.17)
    add_rect(sl, x, Inches(3.85), Inches(3.05), Inches(2.8), fill=LGRAY)
    add_rect(sl, x, Inches(3.85), Inches(3.05), Inches(0.7), fill=BLUE)
    add_text(sl, title, x + Inches(0.1), Inches(3.9), Inches(2.85), Inches(0.6),
             font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, desc, x + Inches(0.1), Inches(4.6), Inches(2.85), Inches(1.0),
             font_size=Pt(11), color=DGRAY, wrap=True)

# ══════════════════════════════════════════════════════════
# SLIDE 3 — Mục Tiêu & Phạm Vi
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Mục Tiêu & Phạm Vi", 3)

# 3 goals
goals = [
    ("01", "Fine-tune YOLOv11-N", "Trên bộ dữ liệu kết hợp\nFishEye8K + VisDrone2019"),
    ("02", "Tích hợp SAHI", "Nâng cao phát hiện đối tượng\nnhỏ (người đi bộ xa tâm ảnh)"),
    ("03", "Xây dựng Flask ITS", "Ứng dụng phân tích giao thông\nvới 5 module thông minh"),
]
for i, (num, title, desc) in enumerate(goals):
    x = Inches(0.3 + i * 4.35)
    add_rect(sl, x, Inches(0.85), Inches(4.1), Inches(2.8), fill=LGRAY)
    add_rect(sl, x, Inches(0.85), Inches(1.0), Inches(2.8), fill=NAVY)
    add_text(sl, num, x + Inches(0.1), Inches(1.15), Inches(0.8), Inches(1.0),
             font_size=Pt(36), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(sl, title, x + Inches(1.1), Inches(0.95), Inches(3.0), Inches(0.6),
             font_size=Pt(13), bold=True, color=WHITE)
    add_text(sl, desc, x + Inches(1.1), Inches(1.55), Inches(3.0), Inches(1.0),
             font_size=Pt(11), color=RGBColor(0xCC,0xCC,0xFF))
    add_text(sl, desc, x + Inches(0.1), Inches(3.75), Inches(3.9), Inches(0.8),
             font_size=Pt(11), color=DGRAY, wrap=True)

# correct desc positions
for i, (num, title, desc) in enumerate(goals):
    x = Inches(0.3 + i * 4.35)
    add_text(sl, desc, x + Inches(0.15), Inches(3.65+0), Inches(3.85), Inches(0.8),
             font_size=Pt(11), color=DGRAY, wrap=True)

# 5 classes table
add_text(sl, "Phạm vi: Phát hiện 5 lớp phương tiện giao thông",
         Inches(0.3), Inches(3.78), Inches(12.6), Inches(0.45),
         font_size=Pt(13), bold=True, color=NAVY)

add_table(sl,
    ["Lớp", "Tỷ lệ trong dataset", "Ghi chú"],
    [
        ["Car (Ô tô)", "~45%", "Phổ biến nhất"],
        ["Bus (Xe buýt)", "~8%", "Đối tượng lớn"],
        ["Truck (Xe tải)", "~12%", "Đối tượng lớn, dài"],
        ["Pedestrian (Người đi bộ)", "~20%", "Thách thức nhất — nhỏ & bị che khuất"],
        ["Motorbike (Xe máy)", "~15%", "Phổ biến ở Việt Nam"],
    ],
    x=Inches(0.3), y=Inches(4.25), w=Inches(12.6), h=Inches(2.8),
    col_widths=[Inches(3.5), Inches(2.5), Inches(6.6)],
    font_size=Pt(11)
)

# ══════════════════════════════════════════════════════════
# SLIDE 4 — Lợi Thế Camera Mắt Cá
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Lợi Thế Camera Mắt Cá", 4)

# quote box
add_rect(sl, Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.7), fill=RGBColor(0xFF,0xF3,0xE0))
add_rect(sl, Inches(0.3), Inches(0.85), Inches(0.08), Inches(0.7), fill=ORANGE)
add_text(sl, '"Một camera fisheye thay cho bốn — toàn bộ chức năng ITS tại một điểm lắp đặt."',
         Inches(0.5), Inches(0.9), Inches(12.4), Inches(0.55),
         font_size=Pt(13), italic=True, color=RGBColor(0x8B,0x45,0x13),
         align=PP_ALIGN.CENTER)

add_table(sl,
    ["Tiêu chí", "Camera thường (x4)", "Camera Fisheye (đề tài)"],
    [
        ["Phạm vi ngã tư", "Cần 4 camera", ("1 camera 180°–220°", True, RGBColor(0x1B,0x5E,0x20))],
        ["Chi phí lắp đặt", "Cao (4× thiết bị, cáp, nguồn)", ("Thấp hơn đáng kể", True, RGBColor(0x1B,0x5E,0x20))],
        ["Đếm lưu lượng đa hướng", "Đồng bộ 4 luồng riêng", ("1 luồng, 4 vạch ảo", True, RGBColor(0x1B,0x5E,0x20))],
        ["Mô hình AI", "YOLO thông thường", ("YOLOv11 fine-tune fisheye", True, NAVY)],
        ["Quản lý hạ tầng", "Phức tạp", ("Đơn giản", True, RGBColor(0x1B,0x5E,0x20))],
    ],
    x=Inches(0.3), y=Inches(1.65), w=Inches(12.7), h=Inches(3.8),
    col_widths=[Inches(3.2), Inches(4.5), Inches(5.0)],
    font_size=Pt(12)
)

# 3 benefit boxes at bottom
benefits = [
    ("💰", "Tiết kiệm chi phí", "Giảm ~75% số lượng\nthiết bị lắp đặt"),
    ("📊", "Phân tích đồng bộ", "Dữ liệu 360° từ\nmột nguồn duy nhất"),
    ("🤖", "AI chuyên biệt", "Fine-tune cho ảnh\nfisheye distorted"),
]
for i, (icon, title, desc) in enumerate(benefits):
    x = Inches(0.3 + i * 4.35)
    add_rect(sl, x, Inches(5.6), Inches(4.1), Inches(1.55), fill=LGRAY)
    add_rect(sl, x, Inches(5.6), Inches(4.1), Inches(0.05), fill=BLUE)
    add_text(sl, f"{title}: {desc.replace(chr(10),' ')}", x + Inches(0.1), Inches(5.68), Inches(3.9), Inches(1.35),
             font_size=Pt(12), bold=False, color=NAVY, wrap=True)

# ══════════════════════════════════════════════════════════
# SLIDE 5 — Cơ Sở Lý Thuyết: Mô Hình Chiếu Fisheye
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Cơ Sở Lý Thuyết: Mô Hình Chiếu Fisheye", 5)

# formula box
add_rect(sl, Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.8), fill=RGBColor(0xE8,0xEF,0xF8))
add_text(sl, "Phương trình chiếu tổng quát:  r(θ) = f · g(θ)",
         Inches(0.5), Inches(0.9), Inches(8.0), Inches(0.5),
         font_size=Pt(14), bold=True, color=NAVY, font_name='Calibri')
add_text(sl, "f = tiêu cự    θ = góc tới    g(θ) = hàm chiếu đặc trưng",
         Inches(0.5), Inches(1.25), Inches(12.0), Inches(0.35),
         font_size=Pt(11), color=DGRAY, italic=True)

add_table(sl,
    ["Mô hình", "Công thức g(θ)", "Ứng dụng thực tế"],
    [
        [("Equidistant ★", True, NAVY), "r = f·θ", "Camera an ninh — phổ biến nhất → Đề tài sử dụng"],
        ["Equisolid", "r = 2f·sin(θ/2)", "Bảo toàn diện tích, đo lường khoa học"],
        ["Orthographic", "r = f·sin(θ)", "Góc nhìn tối đa 180°, thiên văn học"],
        ["Stereographic", "r = 2f·tan(θ/2)", "Bảo toàn góc (conformal projection)"],
        ["Rectilinear", "r = f·tan(θ)", "Không méo đường thẳng, góc < 180°"],
    ],
    x=Inches(0.3), y=Inches(1.75), w=Inches(12.7), h=Inches(3.1),
    col_widths=[Inches(2.8), Inches(3.2), Inches(6.7)],
    font_size=Pt(12)
)

# barrel distortion effects
add_rect(sl, Inches(0.3), Inches(4.95), Inches(12.7), Inches(0.4), fill=NAVY)
add_text(sl, "Ảnh hưởng của Barrel Distortion lên bài toán nhận diện",
         Inches(0.5), Inches(4.97), Inches(12.0), Inches(0.35),
         font_size=Pt(13), bold=True, color=WHITE)

effects = [
    ("① Hình dạng méo", "Bbox chữ nhật kém phù hợp cho đối tượng vùng biên"),
    ("② Feature map lệch", "CNN nhận feature phân phối không đồng đều"),
    ("③ Kích thước phi tuyến", "Cùng đối tượng, khác vị trí → kích thước khác nhau"),
]
for i, (num, desc) in enumerate(effects):
    x = Inches(0.3 + i * 4.27)
    add_rect(sl, x, Inches(5.42), Inches(4.1), Inches(1.75), fill=LGRAY)
    add_text(sl, num, x + Inches(0.1), Inches(5.5), Inches(3.9), Inches(0.45),
             font_size=Pt(12), bold=True, color=NAVY)
    add_text(sl, desc, x + Inches(0.1), Inches(5.98), Inches(3.9), Inches(1.0),
             font_size=Pt(11), color=DGRAY, wrap=True)

# ══════════════════════════════════════════════════════════
# SLIDE 6 — Thuật Toán Biến Đổi Fisheye
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Thuật Toán Biến Đổi Fisheye (fisheye.py)", 6)

# left panel - algorithm
add_rect(sl, Inches(0.3), Inches(0.85), Inches(6.2), Inches(5.8), fill=LGRAY)
add_text(sl, "Hàm to_fisheye() — Inverse Mapping + Bilinear Interpolation",
         Inches(0.4), Inches(0.92), Inches(6.0), Inches(0.5),
         font_size=Pt(13), bold=True, color=NAVY)

steps = [
    ("Bước 1:", "Tính khoảng cách chuẩn hóa r đến tâm ảnh ∈ [0, 1]"),
    ("Bước 2:", "Áp dụng biến đổi bán kính phi tuyến:\n  r' = r^(1 + strength)"),
    ("Bước 3:", "Nội suy song tuyến (bilinear) từ ảnh nguồn\n  NumPy vectorized — không vòng lặp per-pixel"),
]
for i, (num, desc) in enumerate(steps):
    y = Inches(1.5 + i * 1.3)
    add_rect(sl, Inches(0.4), y, Inches(0.8), Inches(0.45), fill=ORANGE)
    add_text(sl, num, Inches(0.4), y, Inches(0.8), Inches(0.45),
             font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, desc, Inches(1.3), y, Inches(5.0), Inches(1.1),
             font_size=Pt(11), color=DGRAY, wrap=True)

add_text(sl, "Tham số tối ưu (thực nghiệm):",
         Inches(0.4), Inches(5.4), Inches(5.0), Inches(0.4),
         font_size=Pt(12), bold=True, color=NAVY)
add_text(sl, "strength = 0.5  ·  radius = 0.85\n→ Tương đồng camera fisheye equidistant thực tế",
         Inches(0.4), Inches(5.82), Inches(5.9), Inches(0.7),
         font_size=Pt(11), color=DGRAY)

# right panel - bbox transform
add_rect(sl, Inches(6.7), Inches(0.85), Inches(6.3), Inches(5.8), fill=RGBColor(0xFF,0xF3,0xE0))
add_text(sl, "Chuyển đổi Bounding Box",
         Inches(6.8), Inches(0.92), Inches(6.1), Inches(0.5),
         font_size=Pt(13), bold=True, color=ORANGE)

bbox_steps = [
    "Lấy mẫu 32 điểm trên toàn bộ chu vi bbox\n(thay vì chỉ 4 góc — chính xác hơn ở vùng biên)",
    "Áp dụng fisheye transform cho từng điểm",
    "Tính axis-aligned bounding box nhỏ nhất\nbao quanh 32 điểm đã biến đổi",
    "Độ chính xác cao hơn đáng kể ở vùng biên\n(nơi gradient biến dạng lớn nhất)",
]
for i, step in enumerate(bbox_steps):
    y = Inches(1.5 + i * 1.1)
    add_rect(sl, Inches(6.8), y, Inches(0.35), Inches(0.35), fill=NAVY)
    add_text(sl, str(i+1), Inches(6.8), y, Inches(0.35), Inches(0.35),
             font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, step, Inches(7.25), y, Inches(5.6), Inches(1.0),
             font_size=Pt(11), color=DGRAY, wrap=True)

# profiles & perf
add_text(sl, "4 Profile: standard · extreme · subtle · traffic_camera",
         Inches(6.8), Inches(5.6), Inches(6.1), Inches(0.4),
         font_size=Pt(11), bold=True, color=NAVY)
add_text(sl, "Hiệu năng: ~50–100ms/frame ảnh 1080p (CPU, vectorized)",
         Inches(6.8), Inches(6.0), Inches(6.1), Inches(0.35),
         font_size=Pt(11), color=DGRAY, italic=True)

# ══════════════════════════════════════════════════════════
# SLIDE 7 — Kiến Trúc YOLOv11
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Kiến Trúc YOLOv11-N", 7)

# top stats
stats = [
    ("2.6M", "Tham số"),
    ("6.5", "GFLOPs"),
    ("~5.3MB", "Kích thước weights"),
    ("Nano", "Phiên bản"),
]
for i, (val, label) in enumerate(stats):
    x = Inches(0.3 + i * 3.18)
    add_rect(sl, x, Inches(0.85), Inches(3.0), Inches(1.1), fill=NAVY)
    add_text(sl, val, x + Inches(0.1), Inches(0.95), Inches(2.8), Inches(0.6),
             font_size=Pt(30), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(sl, label, x + Inches(0.1), Inches(1.52), Inches(2.8), Inches(0.35),
             font_size=Pt(12), color=RGBColor(0xCC,0xCC,0xFF), align=PP_ALIGN.CENTER)

# 3 architecture components
components = [
    ("BACKBONE", "C3k2 — Cross Stage Partial with 2 bottlenecks",
     "• Chia đầu vào 2 nhánh (bottleneck + skip), ghép lại\n• Giảm 3→2 bottleneck block → ít params hơn C3 gốc\n• Cải thiện luồng gradient → hội tụ ổn định khi fine-tune",
     NAVY),
    ("NECK", "FPN + PAN + AIFI Transformer",
     "• FPN: đặc trưng lớp sâu → lớp nông (hỗ trợ đối tượng nhỏ)\n• AIFI: self-attention kế thừa RT-DETR\n  — Mô hình hóa quan hệ dài hạn trong feature map\n  — Giới hạn intra-scale → kiểm soát chi phí tính toán",
     BLUE),
    ("HEAD", "Anchor-free Detection Head",
     "• Dự đoán khoảng cách l, r, t, b từ tâm grid cell\n• Không cần chọn anchor size/ratio\n• Tốt hơn với tỷ lệ bất thường (xe buýt nằm ngang)",
     ORANGE),
]
for i, (label, title, desc, color) in enumerate(components):
    x = Inches(0.3 + i * 4.35)
    add_rect(sl, x, Inches(2.1), Inches(4.1), Inches(4.55), fill=LGRAY)
    add_rect(sl, x, Inches(2.1), Inches(4.1), Inches(0.55), fill=color)
    add_text(sl, label, x + Inches(0.1), Inches(2.15), Inches(3.9), Inches(0.45),
             font_size=Pt(13), bold=True, color=WHITE)
    add_text(sl, title, x + Inches(0.1), Inches(2.72), Inches(3.9), Inches(0.5),
             font_size=Pt(12), bold=True, color=color if color != ORANGE else NAVY)
    add_text(sl, desc, x + Inches(0.1), Inches(3.28), Inches(3.9), Inches(3.1),
             font_size=Pt(10.5), color=DGRAY, wrap=True)

add_text(sl, "So sánh với YOLOv8-N: mAP@0.5:0.95 = 39.5% vs 37.3%  ·  Ít tham số hơn (2.6M vs 3.2M)",
         Inches(0.3), Inches(6.72), Inches(12.7), Inches(0.4),
         font_size=Pt(11), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 8 — Kỹ Thuật SAHI
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Kỹ Thuật SAHI — Sliced Aided Hyper Inference", 8)

# problem statement
add_rect(sl, Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.65), fill=RGBColor(0xFF,0xEB,0xEE))
add_rect(sl, Inches(0.3), Inches(0.85), Inches(0.08), Inches(0.65), fill=RGBColor(0xC6,0x28,0x28))
add_text(sl, "Vấn đề: Người đi bộ xa tâm ảnh fisheye chỉ chiếm 10–30 pixel — model bỏ sót nhiều",
         Inches(0.5), Inches(0.92), Inches(12.4), Inches(0.5),
         font_size=Pt(13), bold=True, color=RGBColor(0xC6,0x28,0x28))

# left: SAHI steps
add_rect(sl, Inches(0.3), Inches(1.6), Inches(5.8), Inches(4.8), fill=LGRAY)
add_text(sl, "Nguyên lý SAHI (Akyon et al., IEEE ICIP 2022)",
         Inches(0.4), Inches(1.65), Inches(5.6), Inches(0.5),
         font_size=Pt(13), bold=True, color=NAVY)
sahi_steps = [
    "① Chia ảnh gốc (640×640) thành các lát nhỏ chồng lấp",
    "② Chạy YOLO inference độc lập trên từng lát",
    "③ Tổng hợp kết quả từ tất cả lát + ảnh gốc",
    "④ NMM (Non-Maximum Merging) loại bỏ detection trùng lặp",
]
for i, step in enumerate(sahi_steps):
    y = Inches(2.25 + i * 0.95)
    add_rect(sl, Inches(0.4), y, Inches(5.6), Inches(0.8), fill=WHITE)
    add_rect(sl, Inches(0.4), y, Inches(0.06), Inches(0.8), fill=BLUE)
    add_text(sl, step, Inches(0.55), y + Inches(0.05), Inches(5.2), Inches(0.7),
             font_size=Pt(12), color=DGRAY, wrap=True)

# right: results
add_rect(sl, Inches(6.3), Inches(1.6), Inches(6.7), Inches(4.8), fill=RGBColor(0xE8,0xF5,0xE9))
add_text(sl, "Kết Quả Thực Nghiệm",
         Inches(6.4), Inches(1.65), Inches(6.5), Inches(0.5),
         font_size=Pt(13), bold=True, color=RGBColor(0x1B,0x5E,0x20))

# big stat
add_rect(sl, Inches(6.4), Inches(2.25), Inches(6.5), Inches(1.3), fill=WHITE)
add_text(sl, "Pedestrian Recall", Inches(6.5), Inches(2.35), Inches(6.3), Inches(0.4),
         font_size=Pt(13), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(sl, "0.42  →  0.75", Inches(6.5), Inches(2.75), Inches(6.3), Inches(0.7),
         font_size=Pt(30), bold=True, color=RGBColor(0x1B,0x5E,0x20), align=PP_ALIGN.CENTER)
add_text(sl, "+78.6% cải thiện", Inches(6.5), Inches(3.42), Inches(6.3), Inches(0.4),
         font_size=Pt(14), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_multiline(sl, [
    ("Lợi ích trong bài toán fisheye:", True, NAVY, Pt(12)),
    ("• Mỗi lát ảnh có biến dạng cục bộ ít hơn", False, DGRAY, Pt(11)),
    ("• Đối tượng nhỏ ở vùng biên xuất hiện tương đối lớn hơn", False, DGRAY, Pt(11)),
    ("", False, DGRAY, Pt(8)),
    ("Đánh đổi tốc độ:", True, NAVY, Pt(12)),
    ("• 1.8–2.5 giây/ảnh (6–8× chậm hơn standard inference)", False, DGRAY, Pt(11)),
    ("• Phù hợp: phân tích offline, snapshot định kỳ 2–3s", False, DGRAY, Pt(11)),
], Inches(6.4), Inches(3.95), Inches(6.5), Inches(2.4))

# ══════════════════════════════════════════════════════════
# SLIDE 9 — Bộ Dữ Liệu
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Bộ Dữ Liệu: FishEye8K + VisDrone2019", 9)

# left panel: FishEye8K
add_rect(sl, Inches(0.3), Inches(0.85), Inches(6.1), Inches(3.5), fill=LGRAY)
add_rect(sl, Inches(0.3), Inches(0.85), Inches(6.1), Inches(0.5), fill=NAVY)
add_text(sl, "FishEye8K — CVPRW 2023, AI City Challenge",
         Inches(0.4), Inches(0.9), Inches(5.9), Inches(0.4),
         font_size=Pt(12), bold=True, color=WHITE)
add_text(sl, "Camera giám sát Sở Cảnh sát Hsinchu, Đài Loan\nGóc overhead, điều kiện đa dạng (ngày/đêm/mưa/sương mù)",
         Inches(0.4), Inches(1.42), Inches(5.9), Inches(0.65),
         font_size=Pt(11), color=DGRAY, italic=True)

add_table(sl,
    ["Split", "Số ảnh", "Số nhãn", "TB/ảnh"],
    [
        ["Train", "4.230", "112.213", "21,2"],
        ["Validation", "1.058", "—", "—"],
        ["Test", "2.712", "—", "—"],
        [("Tổng", True, NAVY), ("8.000", True, NAVY), ("112.213", True, NAVY), "—"],
    ],
    x=Inches(0.35), y=Inches(2.15), w=Inches(6.0), h=Inches(2.1),
    col_widths=[Inches(1.5), Inches(1.5), Inches(1.8), Inches(1.2)],
    font_size=Pt(11)
)

# right panel: VisDrone
add_rect(sl, Inches(6.7), Inches(0.85), Inches(6.3), Inches(3.5), fill=RGBColor(0xFF, 0xF3, 0xE0))
add_rect(sl, Inches(6.7), Inches(0.85), Inches(6.3), Inches(0.5), fill=ORANGE)
add_text(sl, "VisDrone2019 — IEEE/CVF ICCV 2019, Đại học Thiên Tân",
         Inches(6.8), Inches(0.9), Inches(6.1), Inches(0.4),
         font_size=Pt(12), bold=True, color=WHITE)
add_text(sl, "Ảnh UAV, độ cao 10–70m, góc nghiêng + overhead\nÁnh xạ 10 lớp → 5 lớp FishEye8K",
         Inches(6.8), Inches(1.42), Inches(6.1), Inches(0.65),
         font_size=Pt(11), color=DGRAY, italic=True)

add_table(sl,
    ["Split", "Số ảnh", "Số nhãn"],
    [
        ["Train", "6.471", "343.205"],
        ["Val", "548", "38.759"],
        ["Test-Dev", "1.610", "75.102"],
    ],
    x=Inches(6.75), y=Inches(2.15), w=Inches(6.2), h=Inches(2.1),
    col_widths=[Inches(2.0), Inches(2.0), Inches(2.2)],
    font_size=Pt(11)
)

# pipeline arrow
add_rect(sl, Inches(0.3), Inches(4.45), Inches(12.7), Inches(0.45), fill=NAVY)
add_text(sl, "Pipeline chuyển đổi VisDrone → Fisheye",
         Inches(0.5), Inches(4.48), Inches(8.0), Inches(0.38),
         font_size=Pt(13), bold=True, color=WHITE)

pipeline = ["Đọc ảnh +\nnhãn", "to_fisheye()\nstrength=0.5", "32-point bbox\ntransform", "Lọc bbox\nhợp lệ", "Lưu format\nYOLO"]
for i, step in enumerate(pipeline):
    x = Inches(0.3 + i * 2.54)
    add_rect(sl, x, Inches(4.98), Inches(2.35), Inches(1.2), fill=BLUE if i%2==0 else LGRAY)
    add_text(sl, step, x + Inches(0.1), Inches(5.05), Inches(2.15), Inches(1.0),
             font_size=Pt(10.5), bold=True if i%2==0 else False,
             color=WHITE if i%2==0 else NAVY, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(sl, "→", Inches(0.3 + i*2.54 + 2.38), Inches(5.32), Inches(0.25), Inches(0.5),
                 font_size=Pt(16), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# combined dataset
add_rect(sl, Inches(0.3), Inches(6.28), Inches(12.7), Inches(0.9), fill=RGBColor(0xE8,0xEF,0xF8))
add_text(sl, "Dataset kết hợp:  Train 11.296 ảnh · 406.355 nhãn  (TB 35,97 nhãn/ảnh)    |    Val 1.768 ảnh · ~58.000 nhãn",
         Inches(0.5), Inches(6.35), Inches(12.4), Inches(0.5),
         font_size=Pt(13), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 10 — Cấu Hình Huấn Luyện
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Cấu Hình Huấn Luyện", 10)

# env info
add_rect(sl, Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.6), fill=RGBColor(0xE8,0xEF,0xF8))
add_text(sl, "Môi trường: Kaggle Notebooks  ·  GPU Tesla P100-PCIE-16GB  ·  17.1 GB VRAM  ·  RAM 25 GB  ·  CUDA 12.1",
         Inches(0.5), Inches(0.92), Inches(12.4), Inches(0.45),
         font_size=Pt(12), color=NAVY, italic=False, align=PP_ALIGN.CENTER)

add_table(sl,
    ["Siêu tham số", "Cơ Bản (CB)", "Nâng Cao (NC)"],
    [
        ["model", "yolo11n.pt", "yolo11n.pt"],
        ["epochs", "50", ("80", True, RGBColor(0x1B,0x5E,0x20))],
        ["batch_size", "16", "16"],
        ["img_size", "640", ("960", True, RGBColor(0x1B,0x5E,0x20))],
        ["optimizer", "AdamW", ("SGD + Cosine LR", True, RGBColor(0x1B,0x5E,0x20))],
        ["lr0", "0.0005", "0.01"],
        ["freeze backbone", "Không", ("Có (10 lớp đầu)", True, RGBColor(0x1B,0x5E,0x20))],
        ["Dataset", "FishEye8K", ("FishEye8K + VisDrone", True, RGBColor(0x1B,0x5E,0x20))],
        ["SAHI", "Không", ("Có", True, RGBColor(0x1B,0x5E,0x20))],
        ["mosaic / mixup", "1.0 / 0.05", "0.8 / 0.15"],
    ],
    x=Inches(0.3), y=Inches(1.55), w=Inches(8.5), h=Inches(4.95),
    col_widths=[Inches(2.8), Inches(2.6), Inches(3.1)],
    font_size=Pt(11)
)

# right panel - class imbalance techniques
add_rect(sl, Inches(9.0), Inches(1.55), Inches(4.1), Inches(4.95), fill=RGBColor(0xFF, 0xF3, 0xE0))
add_text(sl, "Kỹ thuật cân bằng dữ liệu",
         Inches(9.1), Inches(1.62), Inches(3.9), Inches(0.5),
         font_size=Pt(12), bold=True, color=ORANGE)
techniques = [
    ("Copy-Paste Aug", "Ưu tiên lớp thiểu số\nBus, Truck"),
    ("Class Weight", "Tự động điều chỉnh\ntheo tần suất lớp"),
    ("Oversampling", "Ưu tiên ảnh chứa\nnhiều Bus/Truck"),
]
for i, (t, d) in enumerate(techniques):
    y = Inches(2.2 + i * 1.3)
    add_rect(sl, Inches(9.1), y, Inches(3.8), Inches(1.15), fill=WHITE)
    add_text(sl, t, Inches(9.2), y + Inches(0.08), Inches(3.6), Inches(0.4),
             font_size=Pt(12), bold=True, color=NAVY)
    add_text(sl, d, Inches(9.2), y + Inches(0.52), Inches(3.6), Inches(0.55),
             font_size=Pt(11), color=DGRAY)

add_text(sl, "Thời gian huấn luyện:  CB ~3.8 giờ  ·  NC ~6.8 giờ (img_size=960, 80 epoch)",
         Inches(0.3), Inches(6.65), Inches(12.7), Inches(0.5),
         font_size=Pt(12), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 11 — Kết Quả Huấn Luyện Chi Tiết
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Kết Quả Huấn Luyện — Theo Từng Lớp", 11)

add_text(sl, "CB = Cơ Bản (FishEye8K only)   ·   NC = Nâng Cao (+VisDrone+Freeze+SAHI)",
         Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.4),
         font_size=Pt(11), italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

add_table(sl,
    ["Lớp", "Precision\nCB / NC", "Recall\nCB / NC", "mAP@0.5\nCB / NC", "F1\nCB / NC"],
    [
        ["Car", "0.710 / 0.920", "0.680 / 0.840", "0.720 / 0.910", "0.690 / 0.878"],
        ["Bus", "0.580 / 0.840", "0.520 / 0.700", "0.590 / 0.820", "0.550 / 0.764"],
        ["Truck", "0.600 / 0.850", "0.650 / 0.730", "0.610 / 0.830", "0.620 / 0.785"],
        ["Pedestrian", "0.720 / 0.830", "0.420 / 0.750", "0.550 / 0.850", "0.530 / 0.788"],
        ["Motorbike", "0.640 / 0.906", "0.580 / 0.790", "0.626 / 0.900", "0.610 / 0.845"],
        [("ALL (mean)", True, WHITE), ("0.650 / 0.869", True, WHITE), ("0.570 / 0.762", True, WHITE),
         ("0.619 / 0.862", True, WHITE), ("0.600 / 0.812", True, WHITE)],
    ],
    x=Inches(0.3), y=Inches(1.35), w=Inches(12.7), h=Inches(4.0),
    col_widths=[Inches(2.5), Inches(2.55), Inches(2.55), Inches(2.55), Inches(2.55)],
    font_size=Pt(12)
)

# fix last row bg
# highlight boxes
highlights = [
    ("mAP@0.5 Cơ Bản", "0.619", "baseline"),
    ("mAP@0.5 Nâng Cao", "0.862", "best"),
    ("Cải thiện", "+39.3%", "improvement"),
    ("Pedestrian Recall", "0.42→0.75", "sahi"),
]
for i, (label, val, kind) in enumerate(highlights):
    x = Inches(0.3 + i * 3.18)
    bg = NAVY if kind == "best" else (ORANGE if kind == "improvement" else BLUE if kind == "sahi" else LGRAY)
    fc = WHITE if kind in ("best","improvement","sahi") else NAVY
    add_rect(sl, x, Inches(5.45), Inches(3.1), Inches(1.2), fill=bg)
    add_text(sl, val, x + Inches(0.1), Inches(5.58), Inches(2.9), Inches(0.6),
             font_size=Pt(24), bold=True, color=WHITE if kind != "baseline" else NAVY,
             align=PP_ALIGN.CENTER)
    add_text(sl, label, x + Inches(0.1), Inches(6.18), Inches(2.9), Inches(0.38),
             font_size=Pt(11), color=WHITE if kind != "baseline" else NAVY,
             align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 12 — So Sánh 2 Phiên Bản
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "So Sánh 2 Phiên Bản & Phân Tích", 12)

add_table(sl,
    ["Phiên bản", "mAP@0.5", "mAP@0.5:0.95", "Precision", "Recall", "FPS (GPU)"],
    [
        ["Cơ Bản (FishEye8K)", "0.619", "0.363", "0.650", "0.570", "~41 FPS"],
        [("Nâng Cao (+VisDrone\n+SAHI+Freeze)", True, NAVY),
         ("0.862", True, RGBColor(0x1B,0x5E,0x20)),
         ("0.572", True, RGBColor(0x1B,0x5E,0x20)),
         ("0.869", True, RGBColor(0x1B,0x5E,0x20)),
         ("0.762", True, RGBColor(0x1B,0x5E,0x20)),
         "~12 FPS (SAHI)"],
    ],
    x=Inches(0.3), y=Inches(0.85), w=Inches(12.7), h=Inches(1.5),
    col_widths=[Inches(2.8), Inches(2.0), Inches(2.2), Inches(2.0), Inches(2.0), Inches(1.7)],
    font_size=Pt(11)
)

# improvement boxes
improvements = [
    ("mAP@0.5", "0.619 → 0.862", "+39.3%"),
    ("mAP@0.5:0.95", "0.363 → 0.572", "+57.6%"),
    ("Precision", "0.650 → 0.869", "+33.7%"),
    ("Recall", "0.570 → 0.762", "+33.7%"),
]
for i, (metric, change, pct) in enumerate(improvements):
    x = Inches(0.3 + i * 3.18)
    add_rect(sl, x, Inches(2.5), Inches(3.1), Inches(1.5), fill=LGRAY)
    add_rect(sl, x, Inches(2.5), Inches(3.1), Inches(0.06), fill=RGBColor(0x1B,0x5E,0x20))
    add_text(sl, pct, x + Inches(0.1), Inches(2.58), Inches(2.9), Inches(0.7),
             font_size=Pt(30), bold=True, color=RGBColor(0x1B,0x5E,0x20), align=PP_ALIGN.CENTER)
    add_text(sl, metric, x + Inches(0.1), Inches(3.3), Inches(2.9), Inches(0.38),
             font_size=Pt(11), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(sl, change, x + Inches(0.1), Inches(3.68), Inches(2.9), Inches(0.3),
             font_size=Pt(10), color=DGRAY, align=PP_ALIGN.CENTER)

# analysis
add_rect(sl, Inches(0.3), Inches(4.15), Inches(12.7), Inches(0.45), fill=NAVY)
add_text(sl, "Phân Tích Yếu Tố Cải Thiện",
         Inches(0.5), Inches(4.18), Inches(12.0), Inches(0.38),
         font_size=Pt(13), bold=True, color=WHITE)

analyses = [
    ("Đóng băng backbone\n(freeze=10)", "Bảo toàn đặc trưng pretrained\nCOCO, tránh overfitting"),
    ("Dataset VisDrone", "Tăng gấp đôi số mẫu\n→ cải thiện Bus, Truck, Pedestrian"),
    ("SAHI", "Pedestrian Recall\n0.42 → 0.75 (+78.6%)"),
    ("Đánh đổi tốc độ", "41 FPS → 12 FPS\nChấp nhận được cho offline"),
]
for i, (title, desc) in enumerate(analyses):
    x = Inches(0.3 + i * 3.18)
    add_rect(sl, x, Inches(4.68), Inches(3.1), Inches(2.45), fill=RGBColor(0xE8,0xEF,0xF8))
    add_text(sl, title, x + Inches(0.1), Inches(4.75), Inches(2.9), Inches(0.75),
             font_size=Pt(12), bold=True, color=NAVY)
    add_text(sl, desc, x + Inches(0.1), Inches(5.55), Inches(2.9), Inches(1.35),
             font_size=Pt(11), color=DGRAY, wrap=True)

# ══════════════════════════════════════════════════════════
# SLIDE 13 — Kiến Trúc Hệ Thống
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Kiến Trúc Hệ Thống Ứng Dụng", 13)

# layered architecture diagram
layers = [
    ("Frontend (SPA)", "templates/index.html + static/js/*\nrouter.js · appState.js · api.js", RGBColor(0x42,0x8E,0xC4)),
    ("Flask Application Factory", "create_app(): Config → DB → Blueprints → Extended Routes → Logging → atexit cleanup", NAVY),
    ("Blueprints (routes/)", "core · detect · history · external_camera · examples\nroutes_extended.py: analytics · alerts · incidents · speed · congestion · ALPR · export", BLUE),
    ("Module Nghiệp Vụ (services/)", "video_detect · fisheye · speed_estimator · line_counter\ncongestion_detector · incident_detector · alpr · model_registry · inference", RGBColor(0x1B,0x5E,0x20)),
    ("Dữ Liệu (db.py)", "SQLite / PostgreSQL  +  recent_image_store  (dual DB support via DATABASE_URL)", DGRAY),
]

for i, (label, desc, color) in enumerate(layers):
    y = Inches(0.88 + i * 1.14)
    add_rect(sl, Inches(0.3), y, Inches(12.7), Inches(1.05), fill=color)
    add_text(sl, label, Inches(0.5), y + Inches(0.04), Inches(3.2), Inches(0.45),
             font_size=Pt(12), bold=True, color=WHITE)
    add_text(sl, desc, Inches(3.8), y + Inches(0.04), Inches(9.0), Inches(0.95),
             font_size=Pt(10.5), color=WHITE, wrap=True)
    # arrow between layers
    if i < 4:
        add_text(sl, "▼", Inches(6.5), y + Inches(1.05), Inches(0.5), Inches(0.12),
                 font_size=Pt(10), color=MGRAY, align=PP_ALIGN.CENTER)

# tech stack
add_text(sl, "Công nghệ: Flask 3.x · Python 3.10 · YOLOv11 · OpenCV · EasyOCR · SQLite/PostgreSQL · Docker · GCP",
         Inches(0.3), Inches(6.65), Inches(12.7), Inches(0.4),
         font_size=Pt(11), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 14 — Xử Lý Video Bất Đồng Bộ
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Kiến Trúc Xử Lý Video Bất Đồng Bộ", 14)

# problem
add_rect(sl, Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.62), fill=RGBColor(0xFF,0xEB,0xEE))
add_rect(sl, Inches(0.3), Inches(0.85), Inches(0.08), Inches(0.62), fill=RGBColor(0xC6,0x28,0x28))
add_text(sl, "Vấn đề: Video 30s, 1080p → ~95 giây xử lý (GTX 1060) — Trả về ngay sẽ gây 504 timeout",
         Inches(0.5), Inches(0.93), Inches(12.4), Inches(0.45),
         font_size=Pt(12), bold=True, color=RGBColor(0xC6,0x28,0x28))
add_text(sl, "Giải pháp: Job Queue bất đồng bộ — ThreadPoolExecutor",
         Inches(0.5), Inches(1.38), Inches(12.4), Inches(0.4),
         font_size=Pt(13), bold=True, color=RGBColor(0x1B,0x5E,0x20))

# flow diagram
flow_steps = [
    ("Client POST\n/api/detect", BLUE, "Upload video"),
    ("Validate +\nLưu file tạm", NAVY, "Server side"),
    ("Tạo job_id\n(UUID v4)", NAVY, "DB: pending"),
    ("HTTP 202\nAccepted", RGBColor(0x1B,0x5E,0x20), "< 200ms"),
    ("Client poll\nGET /jobs/{id}", BLUE, "Mỗi 2 giây"),
    ("GET result\nStream MP4", RGBColor(0x1B,0x5E,0x20), "Khi 'done'"),
]
for i, (label, color, sub) in enumerate(flow_steps):
    x = Inches(0.3 + i * 2.17)
    add_rect(sl, x, Inches(1.88), Inches(2.0), Inches(1.05), fill=color)
    add_text(sl, label, x + Inches(0.07), Inches(1.95), Inches(1.86), Inches(0.75),
             font_size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, sub, x + Inches(0.07), Inches(2.72), Inches(1.86), Inches(0.3),
             font_size=Pt(9.5), color=MGRAY, align=PP_ALIGN.CENTER)
    if i < 5:
        add_text(sl, "→", Inches(0.3 + i*2.17 + 2.03), Inches(2.2), Inches(0.22), Inches(0.45),
                 font_size=Pt(14), bold=True, color=NAVY)

# background worker
add_rect(sl, Inches(0.3), Inches(3.2), Inches(12.7), Inches(2.4), fill=LGRAY)
add_rect(sl, Inches(0.3), Inches(3.2), Inches(0.08), Inches(2.4), fill=ORANGE)
add_text(sl, "Background Worker Thread — VideoJobQueue.submit()",
         Inches(0.5), Inches(3.25), Inches(12.0), Inches(0.5),
         font_size=Pt(13), bold=True, color=NAVY)

worker_steps = [
    "Đọc từng frame video",
    "YOLO inference",
    "Speed + Congestion overlay",
    "Ghi annotation → Ghép annotated.mp4",
    "Cập nhật DB: running (% hoàn thành)",
    "Khi xong: DB: done, output_path",
]
for i, step in enumerate(worker_steps):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(3.82 + row * 0.7)
    add_rect(sl, x, y, Inches(0.3), Inches(0.3), fill=ORANGE)
    add_text(sl, str(i+1), x, y, Inches(0.3), Inches(0.3),
             font_size=Pt(9), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, step, x + Inches(0.38), y, Inches(3.75), Inches(0.3),
             font_size=Pt(11), color=DGRAY)

# config
add_rect(sl, Inches(0.3), Inches(5.75), Inches(12.7), Inches(0.65), fill=RGBColor(0xE8,0xEF,0xF8))
config_items = [
    "max_workers = 2",
    "max_queue_size = 10",
    "Job states: pending → running → done / failed",
    "Tránh quá tải VRAM GPU",
]
add_text(sl, "  ·  ".join(config_items),
         Inches(0.5), Inches(5.82), Inches(12.4), Inches(0.5),
         font_size=Pt(12), color=NAVY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 15 — 5 Module ITS
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "5 Module Phân Tích Giao Thông ITS", 15)

modules = [
    ("01", "SpeedEstimator", "IoU tracking liên frame → pixel displacement → km/h\nCảnh báo vi phạm tốc độ với màu sắc trực quan", NAVY),
    ("02", "CongestionDetector", "ROI mật độ → 4 mức Level of Service\nVisualize overlay màu theo mức độ tắc nghẽn", BLUE),
    ("03", "IncidentDetector", "3 loại sự cố: đỗ sai · dừng đột ngột · ngược chiều\nCảnh báo webhook không block pipeline chính", RGBColor(0xC6,0x28,0x28)),
    ("04", "LineCounter", "4 vạch ảo N/S/E/W → đếm xe theo hướng và giờ\nSử dụng cross-product để xác định hướng di chuyển", ORANGE),
    ("05", "ALPR", "YOLO crop → EasyOCR → chuẩn hóa biển số VN\nLưu DB + tra cứu lịch sử phương tiện", RGBColor(0x1B,0x5E,0x20)),
]

for i, (num, title, desc, color) in enumerate(modules):
    y = Inches(0.88 + i * 1.18)
    add_rect(sl, Inches(0.3), y, Inches(12.7), Inches(1.1), fill=LGRAY)
    add_rect(sl, Inches(0.3), y, Inches(1.0), Inches(1.1), fill=color)
    add_text(sl, num, Inches(0.3), y, Inches(1.0), Inches(1.1),
             font_size=Pt(24), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, Inches(1.45), y + Inches(0.08), Inches(3.5), Inches(0.45),
             font_size=Pt(14), bold=True, color=color)
    add_text(sl, desc, Inches(1.45), y + Inches(0.55), Inches(11.4), Inches(0.55),
             font_size=Pt(11), color=DGRAY, wrap=True)

add_text(sl, "Hỗ trợ chung: alert_manager (webhook) · analytics (heatmap + hourly) · cloud_storage (GCS)",
         Inches(0.3), Inches(6.68), Inches(12.7), Inches(0.4),
         font_size=Pt(11), italic=True, color=BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 16 — SpeedEstimator
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Module SpeedEstimator — IoU Tracking", 16)

# left: algorithm
add_rect(sl, Inches(0.3), Inches(0.85), Inches(6.2), Inches(5.8), fill=LGRAY)
add_text(sl, "Thuật Toán IoU Tracking",
         Inches(0.4), Inches(0.92), Inches(6.0), Inches(0.45),
         font_size=Pt(13), bold=True, color=NAVY)
add_text(sl, "(Đơn giản hơn SORT/DeepSORT — không cần model re-ID)",
         Inches(0.4), Inches(1.38), Inches(6.0), Inches(0.35),
         font_size=Pt(10.5), color=DGRAY, italic=True)

# code box
add_rect(sl, Inches(0.4), Inches(1.8), Inches(6.0), Inches(1.8), fill=NAVY)
add_text(sl, 'Frame t:   D_t = {(bbox_i, class_i, conf_i)}\nFrame t+1: Ma trận IoU giữa tất cả cặp\n           Hungarian matching → max Σ IoU (thr=0.3)\n           Unmatched → new track\n           Track mất max_age=5 frame → xóa',
         Inches(0.5), Inches(1.88), Inches(5.8), Inches(1.6),
         font_size=Pt(10), color=RGBColor(0x90,0xFF,0x90), font_name='Consolas')

# speed formula
add_text(sl, "Công thức tốc độ:",
         Inches(0.4), Inches(3.72), Inches(5.0), Inches(0.4),
         font_size=Pt(12), bold=True, color=NAVY)
add_rect(sl, Inches(0.4), Inches(4.15), Inches(6.0), Inches(1.3), fill=WHITE)
add_text(sl, "v (m/s) = √(Δcx² + Δcy²) × (1/pixels_per_meter) × fps\nv (km/h) = v (m/s) × 3.6\npixels_per_meter = 8.0 px/m  (camera 5m, góc 45°)",
         Inches(0.5), Inches(4.25), Inches(5.8), Inches(1.1),
         font_size=Pt(11), color=NAVY, font_name='Consolas')

add_text(sl, "fisheye_correction=True: áp dụng hệ số bù trừ biến dạng hướng kính vùng biên",
         Inches(0.4), Inches(5.55), Inches(6.0), Inches(0.5),
         font_size=Pt(10.5), color=DGRAY, italic=True, wrap=True)

# right: speed color coding
add_rect(sl, Inches(6.7), Inches(0.85), Inches(6.3), Inches(5.8), fill=RGBColor(0xF9,0xFB,0xFF))
add_text(sl, "Mã Màu Cảnh Báo Tốc Độ",
         Inches(6.8), Inches(0.92), Inches(6.1), Inches(0.45),
         font_size=Pt(13), bold=True, color=NAVY)

speed_levels = [
    ("< 40 km/h", RGBColor(0x1B,0x5E,0x20), "Xanh lá", "Bình thường"),
    ("40–70 km/h", ORANGE, "Cam", "Nhanh"),
    ("> 70 km/h", RGBColor(0xC6,0x28,0x28), "Đỏ", "⚠ Vượt tốc"),
]
for i, (speed, color, col_name, label) in enumerate(speed_levels):
    y = Inches(1.55 + i * 1.45)
    add_rect(sl, Inches(6.8), y, Inches(6.1), Inches(1.25), fill=color)
    add_text(sl, speed, Inches(6.95), y + Inches(0.15), Inches(3.0), Inches(0.55),
             font_size=Pt(24), bold=True, color=WHITE)
    add_text(sl, f"{label}  ({col_name})", Inches(6.95), y + Inches(0.72), Inches(5.8), Inches(0.42),
             font_size=Pt(12), color=WHITE)

add_text(sl, "Cảnh báo vi phạm: Tốc độ > ngưỡng trong ≥ 3 frame liên tiếp → bản ghi vi phạm",
         Inches(6.8), Inches(5.9), Inches(6.1), Inches(0.6),
         font_size=Pt(11), bold=True, color=NAVY, wrap=True)

# ══════════════════════════════════════════════════════════
# SLIDE 17 — CongestionDetector
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Module CongestionDetector — Level of Service", 17)

# method
add_rect(sl, Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.65), fill=RGBColor(0xE8,0xEF,0xF8))
add_text(sl, "Phương pháp: Phân tích mật độ trong ROI (tọa độ chuẩn hóa x1,y1,x2,y2 ∈ [0,1]²)",
         Inches(0.5), Inches(0.9), Inches(12.4), Inches(0.55),
         font_size=Pt(12), bold=True, color=NAVY)

# weights
add_text(sl, "Trọng số theo loại xe:  Motorbike 0.5  ·  Car 1.0  ·  Truck 2.0  ·  Bus 2.5  ·  Pedestrian 0.3",
         Inches(0.3), Inches(1.58), Inches(12.7), Inches(0.4),
         font_size=Pt(12), color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# 4 LoS levels
los_levels = [
    ("FREE\nThông thoáng", "density < 0.3", "Luồng giao thông\nbình thường", RGBColor(0x1B,0x5E,0x20)),
    ("MODERATE\nVừa phải", "0.3 ≤ density < 0.6", "Lưu lượng trung bình\nkhông cần can thiệp", RGBColor(0xF5,0x7F,0x17)),
    ("HEAVY\nNặng", "0.6 ≤ density < 0.9", "Tắc nghẽn cục bộ\ncần chú ý", ORANGE),
    ("SEVERE\nNghiêm trọng", "density ≥ 0.9", "Tắc nghẽn nghiêm trọng\ncần can thiệp ngay", RGBColor(0xC6,0x28,0x28)),
]
for i, (level, condition, meaning, color) in enumerate(los_levels):
    x = Inches(0.3 + i * 3.18)
    add_rect(sl, x, Inches(2.08), Inches(3.05), Inches(3.5), fill=color)
    add_text(sl, level, x + Inches(0.1), Inches(2.15), Inches(2.85), Inches(1.0),
             font_size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, x + Inches(0.1), Inches(3.2), Inches(2.85), Inches(0.55), fill=RGBColor(0xFF,0xFF,0xFF))
    add_text(sl, condition, x + Inches(0.1), Inches(3.25), Inches(2.85), Inches(0.45),
             font_size=Pt(12), bold=True, color=color, align=PP_ALIGN.CENTER,
             font_name='Consolas')
    add_text(sl, meaning, x + Inches(0.1), Inches(3.82), Inches(2.85), Inches(1.6),
             font_size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

# visualization & alert
add_rect(sl, Inches(0.3), Inches(5.68), Inches(12.7), Inches(0.55), fill=LGRAY)
add_text(sl, "Visualization: Hình chữ nhật bán trong suốt theo màu  ·  Text: tên ROI, số xe/capacity, %  ·  Dashboard tổng hợp",
         Inches(0.5), Inches(5.73), Inches(12.2), Inches(0.45),
         font_size=Pt(11), color=NAVY)
add_rect(sl, Inches(0.3), Inches(6.28), Inches(12.7), Inches(0.45), fill=RGBColor(0xC6,0x28,0x28))
add_text(sl, "Alert: SEVERE → Gửi webhook (không block pipeline chính)  ·  alert_manager async",
         Inches(0.5), Inches(6.33), Inches(12.2), Inches(0.38),
         font_size=Pt(11), bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════
# SLIDE 18 — Analytics & ALPR
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Module Analytics, Heatmap & ALPR", 18)

# left: analytics
add_rect(sl, Inches(0.3), Inches(0.85), Inches(6.2), Inches(5.8), fill=LGRAY)
add_text(sl, "Heatmap Mật Độ (analytics.py)",
         Inches(0.4), Inches(0.92), Inches(6.0), Inches(0.45),
         font_size=Pt(13), bold=True, color=NAVY)

add_rect(sl, Inches(0.4), Inches(1.45), Inches(6.0), Inches(1.65), fill=NAVY)
add_text(sl, "# Mỗi detection tại (cx, cy):\nheatmap_acc[cy, cx] += 1\n# Render định kỳ:\nGaussian blur (sigma=15) → Normalize [0,255]\n→ cv2.COLORMAP_INFERNO → Blend ảnh gốc (alpha=0.4)",
         Inches(0.5), Inches(1.53), Inches(5.8), Inches(1.5),
         font_size=Pt(10), color=RGBColor(0x90,0xFF,0x90), font_name='Consolas')

add_text(sl, "Dashboard Analytics:",
         Inches(0.4), Inches(3.22), Inches(5.0), Inches(0.4),
         font_size=Pt(12), bold=True, color=NAVY)
dashboard_items = [
    "• Biểu đồ lưu lượng theo giờ",
    "• Phân bố loại phương tiện (pie chart)",
    "• Xác định giờ cao điểm",
    "• TOC — Traffic Operations Center: 8 trang SPA",
    "• Export CSV / JSON",
]
add_multiline(sl, [(item, False, DGRAY, Pt(11)) for item in dashboard_items],
              Inches(0.4), Inches(3.68), Inches(6.0), Inches(2.75))

# right: ALPR
add_rect(sl, Inches(6.7), Inches(0.85), Inches(6.3), Inches(5.8), fill=RGBColor(0xFF, 0xF3, 0xE0))
add_text(sl, "Module ALPR — Nhận Dạng Biển Số Việt Nam",
         Inches(6.8), Inches(0.92), Inches(6.1), Inches(0.5),
         font_size=Pt(13), bold=True, color=ORANGE)

alpr_pipeline = [
    "YOLO bbox → Crop vùng xe (ưu tiên xe lớn nhất)",
    "EasyOCR → raw text",
    'Regex chuẩn hóa biển VN:\n"51F-123.45" hoặc "29X1-2345"',
    "Validate: confidence ≥ 0.30, độ dài 7–10 ký tự",
    "Lưu bảng license_plates, tra cứu lịch sử",
]
for i, step in enumerate(alpr_pipeline):
    y = Inches(1.55 + i * 0.95)
    add_rect(sl, Inches(6.8), y, Inches(0.45), Inches(0.45), fill=ORANGE)
    add_text(sl, str(i+1), Inches(6.8), y, Inches(0.45), Inches(0.45),
             font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, step, Inches(7.38), y, Inches(5.5), Inches(0.85),
             font_size=Pt(11), color=DGRAY, wrap=True)

# ══════════════════════════════════════════════════════════
# SLIDE 19 — Giao Diện Web
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Giao Diện Web — 6 Phân Hệ", 19)

add_text(sl, "Kiến trúc Frontend: Vanilla JavaScript SPA — HTML5 · Bootstrap 5 · ES6 Modules",
         Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.4),
         font_size=Pt(12), italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

add_table(sl,
    ["Phân hệ", "Chức năng chi tiết"],
    [
        [("Dashboard", True, NAVY), "KPIs theo thời gian thực: tổng lượt chạy, số đối tượng, biểu đồ phân phối, system logs"],
        [("Workspace (Inference)", True, NAVY), "Kéo-thả ảnh/video · Đồng bộ (ảnh) · Bất đồng bộ qua job queue (video) · Annotated output"],
        [("Live Streams", True, NAVY), "MJPEG stream fisheye · Song song: camera gốc + detection overlay + tốc độ + ùn tắc"],
        [("Run History", True, NAVY), "Lịch sử phiên nhận diện · Tải ảnh/video annotated + metadata JSON"],
        [("System Logs", True, NAVY), "Terminal view · Server-Sent Events (SSE) · Tìm kiếm, lọc theo level"],
        [("TOC / ALPR", True, NAVY), "Trung tâm điều hành · Nhận dạng & tra cứu biển số · Heatmap · Export CSV/JSON"],
    ],
    x=Inches(0.3), y=Inches(1.32), w=Inches(12.7), h=Inches(4.25),
    col_widths=[Inches(3.0), Inches(9.7)],
    font_size=Pt(11)
)

# UX principles
add_rect(sl, Inches(0.3), Inches(5.65), Inches(12.7), Inches(0.45), fill=NAVY)
add_text(sl, "UX Principles",
         Inches(0.5), Inches(5.68), Inches(3.0), Inches(0.38),
         font_size=Pt(12), bold=True, color=WHITE)

ux = [
    "Responsive: 1920×1080 + tablet 1024×768",
    "Progressive disclosure: chi tiết kỹ thuật ẩn mặc định",
    "Error handling: \"File quá lớn. Vui lòng chọn < 500MB\"",
    "Feedback tức thì: loading spinner, toast notification",
]
for i, item in enumerate(ux):
    x = Inches(0.3 + i * 3.18)
    add_rect(sl, x, Inches(6.17), Inches(3.1), Inches(0.98), fill=LGRAY)
    add_text(sl, item, x + Inches(0.1), Inches(6.22), Inches(2.9), Inches(0.88),
             font_size=Pt(10.5), color=NAVY, wrap=True)

# ══════════════════════════════════════════════════════════
# SLIDE 20 — Kiểm Thử
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Kết Quả Kiểm Thử Chức Năng", 20)

add_table(sl,
    ["ID", "Kịch bản kiểm thử", "Kết quả", "Thời gian"],
    [
        ["TC-01", "Upload ảnh JPEG 1920×1080, phát hiện đối tượng", "✅ PASS", "480ms"],
        ["TC-02", "Upload ảnh > 50MB (giới hạn)", "✅ PASS (Error 413)", "50ms"],
        ["TC-03", "Upload video 30s, 1080p, xử lý bất đồng bộ", "✅ PASS (Job async)", "Tạo: 200ms"],
        ["TC-04", "Polling job status khi đang xử lý", "✅ PASS (status + progress%)", "45ms"],
        ["TC-05", "Download video kết quả đã annotate", "✅ PASS (Stream MP4)", "Streaming"],
        ["TC-06", "SAHI inference trên ảnh đông người đi bộ", "✅ PASS", "2.1s"],
        ["TC-07", "API /api/health check", "✅ PASS", "12ms"],
        ["TC-08", "Gửi webhook khi mật độ SEVERE", "✅ PASS", "~100ms"],
        ["TC-09", "Concurrent 3 video jobs đồng thời", "✅ PASS (Queue đúng thứ tự)", "—"],
    ],
    x=Inches(0.3), y=Inches(0.85), w=Inches(12.7), h=Inches(5.55),
    col_widths=[Inches(1.3), Inches(6.0), Inches(3.2), Inches(2.2)],
    font_size=Pt(11)
)

# summary
add_rect(sl, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.7), fill=RGBColor(0xE8,0xF5,0xE9))
add_rect(sl, Inches(0.3), Inches(6.5), Inches(0.08), Inches(0.7), fill=RGBColor(0x1B,0x5E,0x20))
add_text(sl, "Tổng kết:  9/9 Functional Tests PASS   ·   Unit Tests (pytest): 31/31 PASS",
         Inches(0.5), Inches(6.57), Inches(12.4), Inches(0.55),
         font_size=Pt(15), bold=True, color=RGBColor(0x1B,0x5E,0x20),
         align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# SLIDE 21 — Hiệu Năng
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Đánh Giá Hiệu Năng Tổng Thể", 21)

add_text(sl, "Môi trường deployment thực tế: NVIDIA GTX 1060 6GB",
         Inches(0.3), Inches(0.85), Inches(12.7), Inches(0.4),
         font_size=Pt(12), bold=True, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

add_table(sl,
    ["Tác vụ", "Thời gian", "Ghi chú"],
    [
        ["Phát hiện ảnh đơn (1080p)", "380–520ms  (TB: 450ms)", "✅ Đáp ứng yêu cầu < 500ms"],
        ["Xử lý video 1080p, 25fps, 30s", "~95 giây (3.2× realtime)", "~41 FPS effective (P100 Kaggle)"],
        ["SAHI inference", "1.8–2.5s/ảnh", "Recall người đi bộ 32% → 58%"],
        ["RAM sử dụng", "~1.8 GB (idle) / ~3.2 GB (xử lý)", "—"],
        ["VRAM sử dụng", "~2.1 GB (YOLOv11 FP16)", "—"],
        ["API list-jobs", "< 100ms", "—"],
    ],
    x=Inches(0.3), y=Inches(1.35), w=Inches(12.7), h=Inches(3.5),
    col_widths=[Inches(3.5), Inches(4.0), Inches(5.2)],
    font_size=Pt(11)
)

# non-functional requirements
add_rect(sl, Inches(0.3), Inches(4.95), Inches(12.7), Inches(0.45), fill=NAVY)
add_text(sl, "Yêu Cầu Phi Chức Năng Đạt Được",
         Inches(0.5), Inches(4.98), Inches(12.0), Inches(0.38),
         font_size=Pt(13), bold=True, color=WHITE)

nfr = [
    "✅ Xử lý ảnh ≤ 500ms (GPU)",
    "✅ API submit video < 200ms",
    "✅ Uptime ≥ 99% giờ cao điểm",
    "✅ Dual DB (SQLite/PostgreSQL)",
    "✅ Chrome 90+, Firefox 88+, Edge 90+",
]
for i, item in enumerate(nfr):
    x = Inches(0.3 + i * 2.52)
    add_rect(sl, x, Inches(5.48), Inches(2.45), Inches(1.6), fill=RGBColor(0xE8,0xF5,0xE9))
    add_text(sl, item, x + Inches(0.1), Inches(5.58), Inches(2.3), Inches(1.3),
             font_size=Pt(11), bold=True, color=RGBColor(0x1B,0x5E,0x20), wrap=True)

# ══════════════════════════════════════════════════════════
# SLIDE 22 — Kết Quả Đạt Được
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Kết Quả Đạt Được", 22)

# left: model results
add_rect(sl, Inches(0.3), Inches(0.85), Inches(6.1), Inches(5.8), fill=LGRAY)
add_rect(sl, Inches(0.3), Inches(0.85), Inches(6.1), Inches(0.55), fill=NAVY)
add_text(sl, "Phần 1 — Nghiên Cứu Mô Hình",
         Inches(0.4), Inches(0.9), Inches(5.9), Inches(0.45),
         font_size=Pt(13), bold=True, color=WHITE)

model_results = [
    "✅ Pipeline VisDrone → fisheye với 32-point bbox transform",
    "✅ Dataset kết hợp: 11.296 ảnh train, 406.355 nhãn",
    "✅ YOLOv11-N Cơ bản: mAP@0.5 = 0.619",
    "✅ YOLOv11-N Nâng cao: mAP@0.5 = 0.862 (+39.3%)",
    "✅ SAHI: Pedestrian Recall 0.42 → 0.75 (+78.6%)",
    "✅ Chỉ 2.6M tham số — phù hợp edge deployment",
]
add_multiline(sl, [(r, False, RGBColor(0x1B,0x5E,0x20) if "✅" in r else DGRAY, Pt(12)) for r in model_results],
              Inches(0.4), Inches(1.5), Inches(5.9), Inches(5.0), line_spacing=1.3)

# right: app results
add_rect(sl, Inches(6.7), Inches(0.85), Inches(6.3), Inches(5.8), fill=RGBColor(0xE8,0xEF,0xF8))
add_rect(sl, Inches(6.7), Inches(0.85), Inches(6.3), Inches(0.55), fill=BLUE)
add_text(sl, "Phần 2 — Ứng Dụng Giám Sát",
         Inches(6.8), Inches(0.9), Inches(6.1), Inches(0.45),
         font_size=Pt(13), bold=True, color=WHITE)

app_results = [
    "✅ Flask Application Factory + Blueprint modular",
    "✅ Xử lý video bất đồng bộ — không timeout",
    "✅ 5 module ITS: Speed · Congestion · Incident · LineCounter · ALPR",
    "✅ Dashboard TOC + 8 trang SPA",
    "✅ Dual DB, Cloud Storage, CI/CD → GCP (NVIDIA L4)",
    "✅ 9/9 functional tests, 31/31 unit tests PASS",
]
add_multiline(sl, [(r, False, RGBColor(0x1B,0x5E,0x20) if "✅" in r else DGRAY, Pt(12)) for r in app_results],
              Inches(6.8), Inches(1.5), Inches(6.1), Inches(5.0), line_spacing=1.3)

# ══════════════════════════════════════════════════════════
# SLIDE 23 — Hạn Chế & Hướng Phát Triển
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
content_slide(sl)
nav_bar(sl, "Hạn Chế & Hướng Phát Triển", 23)

# limitations table
add_table(sl,
    ["Hạn chế", "Chi tiết"],
    [
        ["Recall Pedestrian thấp", "Recall 0.42 ở CB vì người đi bộ rất nhỏ; SAHI cải thiện lên 0.75 nhưng cần thêm dữ liệu"],
        ["Calibrate tốc độ thủ công", "pixels_per_meter = 8.0 đặt tay; cần quy trình camera calibration tự động"],
        ["ROI tắc nghẽn cố định", "CongestionDetector cần user cấu hình ROI; kém linh hoạt khi camera dịch chuyển"],
        ["Dữ liệu chưa đặc thù VN", "FishEye8K (Đài Loan) + VisDrone (TQ) — xe máy và phong cách lái VN khác biệt"],
        ["Chưa có edge deployment", "Chạy server tập trung; chưa tối ưu cho Jetson Nano, Raspberry Pi"],
    ],
    x=Inches(0.3), y=Inches(0.88), w=Inches(12.7), h=Inches(2.8),
    col_widths=[Inches(3.0), Inches(9.7)],
    font_size=Pt(11)
)

# future directions
add_rect(sl, Inches(0.3), Inches(3.8), Inches(12.7), Inches(0.45), fill=ORANGE)
add_text(sl, "Hướng Phát Triển",
         Inches(0.5), Inches(3.83), Inches(12.0), Inches(0.38),
         font_size=Pt(13), bold=True, color=WHITE)

future = [
    ("Thu thập dữ liệu VN", "2–3 nút giao Hà Nội\n~5.000 ảnh trong 3 tháng\nMotorbike + Pedestrian"),
    ("Kiến trúc mới", "YOLOv12 · RT-DETR\nSelf-attention toàn cục\nXử lý biến dạng tốt hơn"),
    ("Edge computing", "TensorRT/ONNX Runtime\nJetson Orin Nano\nMục tiêu ≥ 25 FPS"),
    ("Dashboard cải tiến", "React.js + WebSocket\nReal-time map (OpenLayers)\nBáo cáo PDF/Excel tự động"),
]
for i, (title, desc) in enumerate(future):
    x = Inches(0.3 + i * 3.18)
    add_rect(sl, x, Inches(4.32), Inches(3.05), Inches(2.82), fill=LGRAY)
    add_rect(sl, x, Inches(4.32), Inches(3.05), Inches(0.5), fill=NAVY)
    add_text(sl, title, x + Inches(0.1), Inches(4.38), Inches(2.85), Inches(0.4),
             font_size=Pt(12), bold=True, color=WHITE)
    add_text(sl, desc, x + Inches(0.1), Inches(4.9), Inches(2.85), Inches(2.1),
             font_size=Pt(11), color=DGRAY, wrap=True)

# ══════════════════════════════════════════════════════════
# SLIDE 24 — Kết Luận
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)

# dark background
add_rect(sl, 0, 0, W, H, fill=NAVY)
add_rect(sl, 0, H - Inches(0.3), W, Inches(0.3), fill=ORANGE)

# left: model conclusions
add_rect(sl, Inches(0.3), Inches(0.6), Inches(6.2), Inches(5.5), fill=RGBColor(0xFF,0xFF,0xFF))
add_rect(sl, Inches(0.3), Inches(0.6), Inches(6.2), Inches(0.6), fill=BLUE)
add_text(sl, "Về Nghiên Cứu Mô Hình",
         Inches(0.4), Inches(0.65), Inches(6.0), Inches(0.5),
         font_size=Pt(14), bold=True, color=WHITE)

model_conc = [
    "• Pipeline 32-point bbox transform → VisDrone fisheye",
    "• YOLOv11-N: mAP@0.5 = 0.862 (+39.3% vs baseline)",
    "• SAHI: Recall người đi bộ 0.42 → 0.75 (+78.6%)",
    "• 2.6M params, 6.5 GFLOPs — phù hợp triển khai thực tế",
]
add_multiline(sl, [(r, False, NAVY, Pt(12)) for r in model_conc],
              Inches(0.4), Inches(1.3), Inches(6.0), Inches(4.5), line_spacing=1.5)

# right: app conclusions
add_rect(sl, Inches(6.8), Inches(0.6), Inches(6.2), Inches(5.5), fill=RGBColor(0xFF,0xFF,0xFF))
add_rect(sl, Inches(6.8), Inches(0.6), Inches(6.2), Inches(0.6), fill=ORANGE)
add_text(sl, "Về Ứng Dụng & Thực Tiễn",
         Inches(6.9), Inches(0.65), Inches(6.0), Inches(0.5),
         font_size=Pt(14), bold=True, color=WHITE)

app_conc = [
    "• 5 module ITS: Speed, Congestion, Incident, LineCounter, ALPR",
    "• 9/9 functional tests · 31/31 unit tests PASS",
    "• Xử lý ảnh < 500ms — sẵn sàng production",
    "• 1 fisheye = 4 camera thường tại ngã tư đô thị",
]
add_multiline(sl, [(r, False, NAVY, Pt(12)) for r in app_conc],
              Inches(6.9), Inches(1.3), Inches(6.0), Inches(4.5), line_spacing=1.5)

# center strip
add_text(sl, "Đồ án góp phần xây dựng nền tảng ITS Việt Nam ứng dụng AI bản địa",
         Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.5),
         font_size=Pt(15), bold=True, color=ORANGE,
         align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════
out_path = r"C:\Using\DATN\BaoVe_NguyenQuocNam_221220938.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
