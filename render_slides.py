"""Render PPTX slides as JPG thumbnails using python-pptx + PIL."""
import os, zipfile
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

PPTX = r'C:\Using\DATN\BaoVe_NguyenQuocNam_221220938.pptx'
OUT  = r'C:\Using\DATN\qa_slides'
os.makedirs(OUT, exist_ok=True)

# Extract all media from the PPTX zip
media = {}
with zipfile.ZipFile(PPTX, 'r') as z:
    for name in z.namelist():
        if name.startswith('ppt/media/'):
            media[os.path.basename(name)] = z.read(name)

prs = Presentation(PPTX)
W, H = 1280, 720

for idx, slide in enumerate(prs.slides, 1):
    img = Image.new('RGB', (W, H), (240, 242, 246))
    draw = ImageDraw.Draw(img)

    # Draw slide number
    draw.rectangle([0, 0, W, 36], fill=(30, 50, 100))
    draw.text((10, 8), f'SLIDE {idx}/17', fill='white')

    # Extract title text
    title_text = ''
    body_lines = []
    for shape in slide.shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        txt = shape.text_frame.text.strip()
        if not txt:
            continue
        # Guess title by font size or position
        if shape.top < 800000 and shape.width > 5000000:
            title_text = txt[:80]
        else:
            body_lines.append(txt[:100])

    # Draw title
    draw.text((10, 46), title_text, fill=(20, 40, 80))

    # Draw body text lines
    y = 90
    for line in body_lines[:12]:
        if y > 550: break
        for sub in line.split('\n')[:3]:
            draw.text((10, y), sub[:120], fill=(60, 60, 60))
            y += 22
        y += 4

    # Place embedded images
    x_img = W // 2
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            try:
                img_bytes = shape.image.blob
                from io import BytesIO
                pil = Image.open(BytesIO(img_bytes)).convert('RGB')
                # Scale to fit right half
                max_w = W - x_img - 20
                max_h = H - 60
                pil.thumbnail((max_w, max_h), Image.LANCZOS)
                paste_x = W - pil.width - 10
                paste_y = 50 + (H - 50 - pil.height) // 2
                img.paste(pil, (paste_x, max(50, paste_y)))
                break
            except Exception as e:
                pass

    # Border
    draw.rectangle([0, 0, W-1, H-1], outline=(100, 120, 180), width=2)
    img.save(os.path.join(OUT, f'slide_{idx:02d}.jpg'), quality=88)
    print(f'Rendered slide {idx}')

print('All slides rendered to', OUT)
