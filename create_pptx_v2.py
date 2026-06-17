# -*- coding: utf-8 -*-
"""Create BaoVe_NguyenQuocNam_221220938.pptx  — v2 with university background + content images"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── image paths ────────────────────────────────────────────
BG        = r"C:\Using\DATN\ref_media\ppt\media\image1.png"       # university template
YOLO_ARCH = r"C:\Using\DATN\ref_media\ppt\media\image2.png"       # YOLO diagram
FISHEYE1  = r"C:\Using\DATN\fisheye_demo\static\results\79c2c63f167d456ca3248a348d30feb8\fisheye_image.jpg"
FISHEYE2  = r"C:\Using\DATN\fisheye_demo\static\results\152e5ed3a9b04e7fa1caedae38c36711\fisheye_image.jpg"
DETECT1   = r"C:\Using\DATN\fisheye_demo\static\results\d56d2f32d964469695956c6eacc46c0b\annotated.jpg"
SAHI_IMG  = r"C:\Using\DATN\fisheye_demo\static\results\9b96c3060b8d45c7a6ab70e25e45a527\annotated.jpg"
STREET    = r"C:\Using\DATN\fisheye_demo\static\results\5d566030b2b9482dadf89871b79e3abc\annotated.jpg"
OVERVIEW_DAY   = r"C:\Using\DATN\fisheye_demo\static\results\20260512024534-a58732b8\overview_annotated.jpg"
OVERVIEW_NIGHT = r"C:\Using\DATN\fisheye_demo\static\results\20260511153233-2b2ef46d\overview_annotated.jpg"
CAM_ANNOTATED  = r"C:\Using\DATN\fisheye_demo\static\results\20260512024534-a58732b8\camera_3_annotated.jpg"
CAM_ORIGINAL   = r"C:\Using\DATN\fisheye_demo\static\results\20260512024534-a58732b8\camera_3_original.jpg"

# ── constants ──────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# Background purple header occupies ~top 0.9" of image1.png
HDR_H = Inches(0.88)   # height of the purple bar in background
CONTENT_Y = Inches(1.0)  # safe content start Y

# colors
NAVY   = RGBColor(0x28, 0x17, 0x6B)   # match university purple
BLUE   = RGBColor(0x44, 0x72, 0xC4)
ORANGE = RGBColor(0xED, 0x7D, 0x31)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DGRAY  = RGBColor(0x44, 0x54, 0x6A)
MGRAY  = RGBColor(0xBF, 0xBF, 0xBF)
TBLHDR = RGBColor(0x28, 0x17, 0x6B)
TBLALT = RGBColor(0xE8, 0xEF, 0xF8)
GREEN  = RGBColor(0x1B, 0x5E, 0x20)
LGREEN = RGBColor(0xE8, 0xF5, 0xE9)
LORANGE= RGBColor(0xFF, 0xF3, 0xE0)
LRED   = RGBColor(0xFF, 0xEB, 0xEE)
RED    = RGBColor(0xC6, 0x28, 0x28)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ══════════════════════════════════════════════════════════

def add_bg(slide):
    """University background — MUST be called first on each slide"""
    return slide.shapes.add_picture(BG, 0, 0, W, H)

def add_img(slide, path, x, y, w, h):
    return slide.shapes.add_picture(path, x, y, w, h)

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=Pt(24), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, font='Calibri', italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return txb

def add_lines(slide, lines, x, y, w, h,
              size=Pt(19), bold=False, color=DGRAY,
              align=PP_ALIGN.LEFT, font='Calibri', spacing=None):
    """lines: list of str OR (text, bold, color, size, italic)"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.05)
    first = True
    for item in lines:
        if isinstance(item, str):
            t, b, c, sz, it = item, bold, color, size, False
        else:
            t  = item[0]
            b  = item[1] if len(item) > 1 else bold
            c  = item[2] if len(item) > 2 else color
            sz = item[3] if len(item) > 3 else size
            it = item[4] if len(item) > 4 else False
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if spacing:
            p.line_spacing = spacing
        r = p.add_run()
        r.text = t
        r.font.size = sz
        r.font.bold = b
        r.font.italic = it
        r.font.name = font
        r.font.color.rgb = c
    return txb

def slide_title(slide, text, num, total=24):
    """White title text over the purple background header bar"""
    add_text(slide, text,
             Inches(0.28), Inches(0.1), Inches(11.5), Inches(0.72),
             size=Pt(36), bold=True, color=WHITE, font='Calibri Light')
    add_text(slide, f"{num}/{total}",
             Inches(12.1), Inches(0.15), Inches(1.0), Inches(0.58),
             size=Pt(19), color=RGBColor(0xDD, 0xDD, 0xFF), align=PP_ALIGN.RIGHT)

def orange_bar(slide, text=""):
    """Thin orange bar at very bottom"""
    add_rect(slide, 0, Inches(7.25), W, Inches(0.25), fill=ORANGE)
    if text:
        add_text(slide, text, Inches(0.3), Inches(7.26), Inches(12.7), Inches(0.22),
                 size=Pt(11), color=WHITE, align=PP_ALIGN.CENTER)

def add_table(slide, headers, rows, x, y, w, h=None,
              col_widths=None, fs=Pt(13), hfs=Pt(14)):
    ncols = len(headers)
    nrows = len(rows) + 1
    if h is None:
        h = Inches(0.36) * nrows + Inches(0.05)
    tbl = slide.shapes.add_table(nrows, ncols, x, y, w, h).table
    tbl.first_row = True
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    def cell_set(cell, text, bold=False, bg=None, fg=WHITE, sz=Pt(13), al=PP_ALIGN.LEFT, it=False):
        tf = cell.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(0.06)
        tf.margin_top = tf.margin_bottom = Inches(0.03)
        p = tf.paragraphs[0]
        p.alignment = al
        r = p.add_run()
        r.text = text
        r.font.size = sz
        r.font.bold = bold
        r.font.italic = it
        r.font.name = 'Calibri'
        r.font.color.rgb = fg
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    for j, h_txt in enumerate(headers):
        cell_set(tbl.cell(0, j), h_txt, bold=True, bg=TBLHDR, fg=WHITE, sz=hfs, al=PP_ALIGN.CENTER)

    for i, row in enumerate(rows):
        bg = TBLALT if i % 2 == 0 else WHITE
        is_last = (i == len(rows) - 1)
        for j, val in enumerate(row):
            b, fg, it = False, DGRAY, False
            if isinstance(val, tuple):
                txt = val[0]
                b   = val[1] if len(val) > 1 else False
                fg  = val[2] if len(val) > 2 else DGRAY
            else:
                txt = str(val)
            if '✅' in txt:
                fg = GREEN; b = True
            # last row of table gets navy bg for total rows
            if is_last and j == 0 and any(kw in txt for kw in ['ALL', 'Tổng', 'mean']):
                bg = TBLHDR; fg = WHITE; b = True
            cell_set(tbl.cell(i+1, j), txt, bold=b, bg=bg, fg=fg, sz=fs)
    return tbl

def img_border(slide, x, y, w, h, path, caption=""):
    """Add image maintaining aspect ratio, centered inside box, navy border"""
    from PIL import Image as _PILImage
    iw, ih = _PILImage.open(path).size
    ratio = iw / ih
    # fit inside box keeping ratio
    fw, fh = w, w / ratio
    if fh > h:
        fh = h; fw = h * ratio
    ox = x + (w - fw) / 2
    oy = y + (h - fh) / 2
    add_rect(slide, ox - Inches(0.035), oy - Inches(0.035),
             fw + Inches(0.07), fh + Inches(0.07), fill=NAVY)
    add_img(slide, path, ox, oy, fw, fh)
    if caption:
        cap_y = oy + fh + Inches(0.04)
        add_text(slide, caption, x, cap_y, w, Inches(0.28),
                 size=Pt(11), color=DGRAY, align=PP_ALIGN.CENTER, italic=True)

# ══════════════════════════════════════════════════════════
# SLIDE 1 — Trang Bìa
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)

# semi-transparent white card in center-left
add_rect(sl, Inches(0.28), Inches(1.05), Inches(8.5), Inches(6.1),
         fill=WHITE, line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))

# university name (replicate reference)
add_text(sl, "TRƯỜNG ĐẠI HỌC GIAO THÔNG VẬN TẢI",
         Inches(0.4), Inches(0.1), Inches(11.0), Inches(0.7),
         size=Pt(24), bold=True, color=WHITE,
         align=PP_ALIGN.LEFT, font='Calibri')

# orange accent line under title box
add_rect(sl, Inches(0.28), Inches(1.05), Inches(8.5), Inches(0.07), fill=ORANGE)

# faculty label
add_text(sl, "KHOA CÔNG NGHỆ THÔNG TIN  ·  ĐỒ ÁN TỐT NGHIỆP",
         Inches(0.45), Inches(1.18), Inches(8.2), Inches(0.42),
         size=Pt(19), color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# main title
add_text(sl, "XÂY DỰNG HỆ THỐNG NHẬN DIỆN\nVẬT THỂ QUA CAMERA MẮT CÁ",
         Inches(0.4), Inches(1.7), Inches(8.3), Inches(1.9),
         size=Pt(38), bold=True, color=NAVY,
         align=PP_ALIGN.CENTER, font='Calibri Light')

# divider
add_rect(sl, Inches(1.0), Inches(3.65), Inches(7.0), Inches(0.03), fill=ORANGE)

# info block
info = [
    ("Sinh viên thực hiện: Nguyễn Quốc Nam", True, NAVY, Pt(17)),
    ("Mã sinh viên: 221220938  ·  Lớp: CNTT1-K63  ·  Khóa 63", False, DGRAY, Pt(19)),
    ("Ngành: Công nghệ Thông tin — Hệ chính quy", False, DGRAY, Pt(19)),
    ("", False, DGRAY, Pt(11)),
    ("Giảng viên hướng dẫn: TS. Nguyễn Đức Dư", True, NAVY, Pt(17)),
    ("", False, DGRAY, Pt(11)),
    ("Hà Nội, 2026", False, DGRAY, Pt(19)),
]
add_lines(sl, info, Inches(0.5), Inches(3.75), Inches(8.0), Inches(2.8),
          align=PP_ALIGN.CENTER, spacing=1.4)

# right panel: 2 landscape images stacked (1.78 ratio)
# cam3_ann: w=4.07" h=2.29"
img_border(sl, Inches(9.05), Inches(1.15), Inches(4.07), Inches(2.29), CAM_ORIGINAL,
           "Camera thực tế — góc nhìn đường phố Việt Nam")
img_border(sl, Inches(9.05), Inches(3.6), Inches(4.07), Inches(2.29), CAM_ANNOTATED,
           "YOLOv11 nhận diện: Car · Pedestrian · Motorbike")

orange_bar(sl, "Đại học Giao thông Vận tải  ·  Hà Nội  ·  2026")

# ══════════════════════════════════════════════════════════
# SLIDE 2 — Đặt Vấn Đề
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Đặt Vấn Đề", 2)

# left column — text (60%)
add_rect(sl, Inches(0.28), CONTENT_Y, Inches(7.7), Inches(2.3), fill=TBLALT)
add_text(sl, "Thực Trạng Giao Thông Việt Nam",
         Inches(0.38), CONTENT_Y + Inches(0.06), Inches(7.5), Inches(0.42),
         size=Pt(17), bold=True, color=NAVY)
add_lines(sl, [
    ("• 7,8 triệu ô tô & 73 triệu xe máy đang lưu hành (Tổng cục Thống kê, 2024)", False, DGRAY, Pt(19)),
    ("• Năm 2023: hơn 10.000 vụ tai nạn giao thông nghiêm trọng", False, DGRAY, Pt(19)),
    ("• Camera CCTV truyền thống: giám sát thủ công — không phân tích tự động", False, DGRAY, Pt(19)),
], Inches(0.38), CONTENT_Y + Inches(0.52), Inches(7.5), Inches(1.65))

add_rect(sl, Inches(0.28), CONTENT_Y + Inches(2.42), Inches(7.7), Inches(2.05), fill=LORANGE)
add_text(sl, "Tại Sao Camera Fisheye?",
         Inches(0.38), CONTENT_Y + Inches(2.5), Inches(7.5), Inches(0.42),
         size=Pt(17), bold=True, color=ORANGE)
add_lines(sl, [
    ("• Góc nhìn 180°–220° — một camera bao phủ toàn bộ ngã tư", True, RGBColor(0x8B, 0x45, 0x13), Pt(19)),
    ("• Tiết kiệm chi phí: thay thế 4 camera thường bằng 1 camera fisheye", False, DGRAY, Pt(19)),
    ("• Nhưng gây barrel distortion → AI thông thường không áp dụng trực tiếp", False, DGRAY, Pt(19)),
], Inches(0.38), CONTENT_Y + Inches(2.95), Inches(7.5), Inches(1.4))

# right panel: 2 landscape images (ratio 1.73/1.78)
# overview top (1.73): w=4.85" h=2.80"
img_border(sl, Inches(8.2), CONTENT_Y, Inches(4.85), Inches(2.8), OVERVIEW_DAY,
           "4-camera fisheye giám sát ngã tư — ban ngày")
img_border(sl, Inches(8.2), CONTENT_Y + Inches(3.0), Inches(4.85), Inches(2.72), OVERVIEW_NIGHT,
           "4-camera fisheye giám sát ngã tư — ban đêm")

# 4 challenge boxes
challenges = [
    ("Bbox kém\nphù hợp", "Đối tượng méo ở vùng biên"),
    ("10–30 px\nnhỏ", "Người đi bộ xa tâm ảnh"),
    ("Phi tuyến\nvị trí", "Kích thước biến đổi theo vị trí"),
    ("Dữ liệu\nkhan hiếm", "Ảnh fisheye nhãn chất lượng cao"),
]
for i, (ttl, dsc) in enumerate(challenges):
    x = Inches(0.28 + i * 3.28)
    add_rect(sl, x, CONTENT_Y + Inches(4.55), Inches(3.1), Inches(1.62), fill=LGRAY)
    add_rect(sl, x, CONTENT_Y + Inches(4.55), Inches(3.1), Inches(0.06), fill=BLUE)
    add_text(sl, ttl, x + Inches(0.1), CONTENT_Y + Inches(4.63), Inches(2.9), Inches(0.55),
             size=Pt(17), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(sl, dsc, x + Inches(0.1), CONTENT_Y + Inches(5.22), Inches(2.9), Inches(0.82),
             size=Pt(14), color=DGRAY, wrap=True)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 3 — Mục Tiêu & Phạm Vi
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Mục Tiêu & Phạm Vi Nghiên Cứu", 3)

goals = [
    ("01", "Fine-tune YOLOv11-N", "Trên bộ dữ liệu kết hợp\nFishEye8K + VisDrone2019", NAVY),
    ("02", "Tích hợp SAHI", "Nâng cao phát hiện đối tượng\nnhỏ (người đi bộ xa tâm ảnh)", BLUE),
    ("03", "Xây dựng Flask ITS", "Ứng dụng phân tích giao thông\nvới 5 module thông minh", ORANGE),
]
for i, (num, ttl, dsc, col) in enumerate(goals):
    x = Inches(0.28 + i * 4.35)
    add_rect(sl, x, CONTENT_Y, Inches(4.1), Inches(2.55), fill=col)
    add_text(sl, num,  x + Inches(0.12), CONTENT_Y + Inches(0.1), Inches(1.0), Inches(1.1),
             size=Pt(44), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, ttl,  x + Inches(1.2), CONTENT_Y + Inches(0.15), Inches(2.8), Inches(0.65),
             size=Pt(24), bold=True, color=WHITE)
    add_text(sl, dsc,  x + Inches(1.2), CONTENT_Y + Inches(0.85), Inches(2.8), Inches(1.5),
             size=Pt(19), color=RGBColor(0xDD, 0xDD, 0xFF), wrap=True)

add_text(sl, "Phạm vi: Phát hiện 5 lớp phương tiện giao thông",
         Inches(0.28), CONTENT_Y + Inches(2.72), Inches(12.7), Inches(0.42),
         size=Pt(17), bold=True, color=NAVY)

add_table(sl,
    ["Lớp đối tượng", "Tỷ lệ trong dataset", "Đặc điểm / Thách thức"],
    [
        ["Car (Ô tô)",              "~45%", "Phổ biến nhất, kích thước trung bình"],
        ["Bus (Xe buýt)",           "~8%",  "Đối tượng lớn, dễ phát hiện"],
        ["Truck (Xe tải)",          "~12%", "Tỷ lệ dài/rộng bất thường"],
        ["Pedestrian (Người đi bộ)","~20%", "Thách thức nhất — nhỏ & bị che khuất"],
        ["Motorbike (Xe máy)",      "~15%", "Phổ biến ở Việt Nam, nhỏ"],
    ],
    x=Inches(0.28), y=CONTENT_Y + Inches(3.2), w=Inches(12.77),
    col_widths=[Inches(3.5), Inches(2.2), Inches(7.07)], fs=Pt(13))

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 4 — Lợi Thế Camera Mắt Cá
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Lợi Thế Camera Mắt Cá", 4)

add_rect(sl, Inches(0.28), CONTENT_Y, Inches(12.77), Inches(0.62), fill=LORANGE)
add_rect(sl, Inches(0.28), CONTENT_Y, Inches(0.07), Inches(0.62), fill=ORANGE)
add_text(sl, '"Một camera fisheye thay cho bốn — toàn bộ chức năng ITS tại một điểm lắp đặt."',
         Inches(0.45), CONTENT_Y + Inches(0.08), Inches(12.4), Inches(0.48),
         size=Pt(17), italic=True, color=RGBColor(0x8B, 0x45, 0x13), align=PP_ALIGN.CENTER)

add_table(sl,
    ["Tiêu chí", "Camera thường (×4)", "Camera Fisheye (đề tài)"],
    [
        ["Phạm vi ngã tư",           "Cần 4 camera, 4 điểm lắp đặt",  ("1 camera, góc 180°–220°", True, GREEN)],
        ["Chi phí lắp đặt",          "Cao (4× thiết bị, cáp, nguồn)", ("Thấp hơn ~75%", True, GREEN)],
        ["Đếm lưu lượng đa hướng",   "Đồng bộ 4 luồng riêng biệt",   ("1 luồng, 4 vạch ảo N/S/E/W", True, GREEN)],
        ["Mô hình AI phát hiện",     "YOLO thông thường",             ("YOLOv11 fine-tune fisheye", True, NAVY)],
        ["Quản lý hạ tầng",          "Phức tạp (cáp, nguồn, đồng bộ)", ("Đơn giản — 1 đầu vào", True, GREEN)],
    ],
    x=Inches(0.28), y=CONTENT_Y + Inches(0.72), w=Inches(12.77), h=Inches(3.0),
    col_widths=[Inches(3.0), Inches(4.5), Inches(5.27)], fs=Pt(13))

# 3 benefit cards
benefits = [
    ("💰 Tiết kiệm chi phí",    "Giảm ~75% số lượng\nthiết bị lắp đặt",           NAVY),
    ("📊 Phân tích đồng bộ",   "Dữ liệu 360° từ\nmột nguồn duy nhất",             BLUE),
    ("🤖 AI chuyên biệt",       "Fine-tune riêng cho\nảnh fisheye distorted",       ORANGE),
    ("⚡ Triển khai nhanh",    "Chỉ cần 1 cáp nguồn\nvà 1 cáp mạng",              RGBColor(0x1B, 0x5E, 0x20)),
]
for i, (ttl, dsc, col) in enumerate(benefits):
    x = Inches(0.28 + i * 3.18)
    add_rect(sl, x, CONTENT_Y + Inches(3.85), Inches(3.07), Inches(1.62), fill=col)
    add_text(sl, ttl, x + Inches(0.12), CONTENT_Y + Inches(3.93), Inches(2.85), Inches(0.48),
             size=Pt(19), bold=True, color=WHITE)
    add_text(sl, dsc, x + Inches(0.12), CONTENT_Y + Inches(4.45), Inches(2.85), Inches(0.95),
             size=Pt(19), color=RGBColor(0xDD, 0xDD, 0xFF), wrap=True)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 5 — Cơ Sở Lý Thuyết: Mô Hình Chiếu Fisheye
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Cơ Sở Lý Thuyết: Mô Hình Chiếu Fisheye", 5)

add_rect(sl, Inches(0.28), CONTENT_Y, Inches(12.77), Inches(0.7), fill=TBLALT)
add_text(sl, "Phương trình chiếu tổng quát:  r(θ) = f · g(θ)   — f: tiêu cự  ·  θ: góc tới  ·  g(θ): hàm đặc trưng từng mô hình",
         Inches(0.42), CONTENT_Y + Inches(0.12), Inches(12.4), Inches(0.5),
         size=Pt(17), bold=True, color=NAVY, align=PP_ALIGN.CENTER, font='Calibri')

add_table(sl,
    ["Mô hình", "Công thức g(θ)", "Đặc điểm & Ứng dụng"],
    [
        [("★ Equidistant", True, NAVY), "r = f·θ",           "Camera an ninh & giao thông — phổ biến nhất → Đề tài sử dụng"],
        ["Equisolid",                  "r = 2f·sin(θ/2)",  "Bảo toàn diện tích, đo lường khoa học"],
        ["Orthographic",               "r = f·sin(θ)",     "Góc nhìn tối đa 180°, thiên văn học"],
        ["Stereographic",              "r = 2f·tan(θ/2)", "Bảo toàn góc cục bộ (conformal projection)"],
        ["Rectilinear",                "r = f·tan(θ)",     "Không méo đường thẳng, góc hẹp < 180°"],
    ],
    x=Inches(0.28), y=CONTENT_Y + Inches(0.78), w=Inches(12.77), h=Inches(2.85),
    col_widths=[Inches(2.6), Inches(3.2), Inches(6.97)], fs=Pt(13))

# barrel distortion section
add_rect(sl, Inches(0.28), CONTENT_Y + Inches(3.74), Inches(12.77), Inches(0.48), fill=NAVY)
add_text(sl, "Ảnh hưởng của Barrel Distortion lên bài toán nhận diện",
         Inches(0.42), CONTENT_Y + Inches(3.79), Inches(12.0), Inches(0.38),
         size=Pt(17), bold=True, color=WHITE)

effects = [
    ("① Hình dạng méo",         "Bbox hình chữ nhật kém\nphù hợp cho đối tượng\nvùng biên ảnh"),
    ("② Feature map lệch",      "CNN xử lý feature map\nphân phối không đều\ntheo tâm/biên ảnh"),
    ("③ Kích thước phi tuyến",  "Cùng đối tượng, khác\nvị trí → kích thước\nkhác nhau đáng kể"),
    ("④ Domain gap",            "Model train ảnh thẳng\nkhông áp dụng trực\ntiếp cho fisheye"),
]
for i, (ttl, dsc) in enumerate(effects):
    x = Inches(0.28 + i * 3.18)
    add_rect(sl, x, CONTENT_Y + Inches(4.15), Inches(3.07), Inches(1.95), fill=LGRAY)
    add_text(sl, ttl, x + Inches(0.1), CONTENT_Y + Inches(4.23), Inches(2.87), Inches(0.42),
             size=Pt(19), bold=True, color=NAVY)
    add_text(sl, dsc, x + Inches(0.1), CONTENT_Y + Inches(4.7), Inches(2.87), Inches(1.25),
             size=Pt(19), color=DGRAY, wrap=True)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 6 — Thuật Toán Biến Đổi Fisheye
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Thuật Toán Biến Đổi Fisheye — fisheye.py", 6)

# left panel — algorithm
add_rect(sl, Inches(0.28), CONTENT_Y, Inches(6.05), Inches(5.85), fill=LGRAY)
add_text(sl, "Hàm to_fisheye() — Inverse Mapping + Bilinear",
         Inches(0.38), CONTENT_Y + Inches(0.08), Inches(5.85), Inches(0.45),
         size=Pt(17), bold=True, color=NAVY)
add_text(sl, "(NumPy vectorized — không vòng lặp per-pixel)",
         Inches(0.38), CONTENT_Y + Inches(0.54), Inches(5.85), Inches(0.32),
         size=Pt(17), color=DGRAY, italic=True)

steps = [
    ("Bước 1", "Tính khoảng cách chuẩn hóa r đến tâm ảnh ∈ [0, 1]"),
    ("Bước 2", "Biến đổi bán kính phi tuyến:\nr' = r^(1 + strength)"),
    ("Bước 3", "Nội suy song tuyến (bilinear) lấy màu từ ảnh nguồn"),
]
for i, (lbl, dsc) in enumerate(steps):
    y = CONTENT_Y + Inches(0.95 + i * 1.28)
    add_rect(sl, Inches(0.38), y, Inches(1.0), Inches(0.42), fill=ORANGE)
    add_text(sl, lbl, Inches(0.38), y, Inches(1.0), Inches(0.42),
             size=Pt(24), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, dsc, Inches(1.48), y, Inches(4.75), Inches(1.1),
             size=Pt(24), color=DGRAY, wrap=True)

add_text(sl, "Tham số: strength=0.5  ·  radius=0.85",
         Inches(0.38), CONTENT_Y + Inches(4.0), Inches(5.85), Inches(0.38),
         size=Pt(19), bold=True, color=NAVY)
add_text(sl, "→ Tương đồng camera equidistant thực tế",
         Inches(0.38), CONTENT_Y + Inches(4.42), Inches(5.85), Inches(0.38),
         size=Pt(17), color=DGRAY, italic=True)
add_text(sl, "4 profile: standard · extreme · subtle · traffic_camera",
         Inches(0.38), CONTENT_Y + Inches(4.83), Inches(5.85), Inches(0.38),
         size=Pt(17), bold=True, color=BLUE)
add_text(sl, "Hiệu năng: ~50–100ms/frame ảnh 1080p (CPU)",
         Inches(0.38), CONTENT_Y + Inches(5.24), Inches(5.85), Inches(0.38),
         size=Pt(17), color=DGRAY, italic=True)

# right panel — bbox transform + images
add_rect(sl, Inches(6.5), CONTENT_Y, Inches(6.55), Inches(2.85), fill=LORANGE)
add_text(sl, "Chuyển đổi Bounding Box — 32 điểm",
         Inches(6.62), CONTENT_Y + Inches(0.08), Inches(6.35), Inches(0.42),
         size=Pt(17), bold=True, color=ORANGE)
bbox_steps = [
    "① Lấy mẫu 32 điểm trên toàn bộ chu vi bbox (thay vì 4 góc)",
    "② Áp dụng fisheye transform cho từng điểm trong 32 điểm",
    "③ Tính axis-aligned bbox bao quanh 32 điểm đã biến đổi",
    "④ Độ chính xác cao hơn đáng kể tại vùng biên ảnh fisheye",
]
for i, s in enumerate(bbox_steps):
    add_text(sl, s, Inches(6.62), CONTENT_Y + Inches(0.58 + i * 0.6), Inches(6.35), Inches(0.55),
             size=Pt(17), color=DGRAY, wrap=True)

# before / after landscape images (ratio 1.78), side by side
# each: w=3.15" h=3.15/1.78=1.77"
img_border(sl, Inches(6.5),  CONTENT_Y + Inches(3.05), Inches(3.15), Inches(1.77), CAM_ORIGINAL,
           "Ảnh gốc — camera thực tế")
img_border(sl, Inches(9.82), CONTENT_Y + Inches(3.05), Inches(3.15), Inches(1.77), CAM_ANNOTATED,
           "Sau detection + bbox transform")

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 7 — Kiến Trúc YOLOv11
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Kiến Trúc YOLOv11-N", 7)

# stats top
stats = [("2.6M", "Tham số"), ("6.5", "GFLOPs"), ("~5.3MB", "Weights"), ("Nano", "Phiên bản")]
for i, (v, l) in enumerate(stats):
    x = Inches(0.28 + i * 3.18)
    add_rect(sl, x, CONTENT_Y, Inches(3.07), Inches(0.98), fill=NAVY)
    add_text(sl, v, x + Inches(0.1), CONTENT_Y + Inches(0.05), Inches(2.87), Inches(0.58),
             size=Pt(36), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(sl, l, x + Inches(0.1), CONTENT_Y + Inches(0.62), Inches(2.87), Inches(0.35),
             size=Pt(17), color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

# YOLO architecture diagram
img_border(sl, Inches(0.28), CONTENT_Y + Inches(1.07), Inches(12.77), Inches(2.55), YOLO_ARCH,
           "Kiến trúc YOLOv11: Backbone (C3k2) → Neck (FPN+AIFI) → Detection Head (Anchor-free)")

# 3 component boxes
comps = [
    ("BACKBONE\nC3k2", "Kế thừa CSP: 2 nhánh (bottleneck + skip)\nGiảm 3→2 block: ít params, gradient tốt hơn\nHội tụ ổn định khi fine-tune dataset nhỏ", NAVY),
    ("NECK\nFPN + AIFI", "FPN: đặc trưng sâu → lớp nông (đối tượng nhỏ)\nAIFI Transformer: self-attention intra-scale\nMô hình hóa quan hệ dài hạn trong feature map", BLUE),
    ("HEAD\nAnchor-free", "Dự đoán l, r, t, b từ tâm grid cell\nKhông cần chọn anchor size/ratio\nTốt hơn với tỷ lệ bất thường (bus nằm ngang)", ORANGE),
]
for i, (ttl, dsc, col) in enumerate(comps):
    x = Inches(0.28 + i * 4.35)
    add_rect(sl, x, CONTENT_Y + Inches(3.72), Inches(4.1), Inches(2.65), fill=LGRAY)
    add_rect(sl, x, CONTENT_Y + Inches(3.72), Inches(4.1), Inches(0.62), fill=col)
    add_text(sl, ttl, x + Inches(0.1), CONTENT_Y + Inches(3.8), Inches(3.9), Inches(0.60),
             size=Pt(17), bold=True, color=WHITE)
    add_text(sl, dsc, x + Inches(0.1), CONTENT_Y + Inches(4.4), Inches(3.9), Inches(1.85),
             size=Pt(17), color=DGRAY, wrap=True)

orange_bar(sl, "So sánh với YOLOv8-N: mAP@0.5:0.95 = 39.5% vs 37.3%  ·  Ít tham số hơn (2.6M vs 3.2M)")

# ══════════════════════════════════════════════════════════
# SLIDE 8 — Kỹ Thuật SAHI
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Kỹ Thuật SAHI — Sliced Aided Hyper Inference", 8)

# problem bar
add_rect(sl, Inches(0.28), CONTENT_Y, Inches(12.77), Inches(0.58), fill=LRED)
add_rect(sl, Inches(0.28), CONTENT_Y, Inches(0.07), Inches(0.58), fill=RED)
add_text(sl, "Vấn đề: Người đi bộ xa tâm ảnh fisheye chỉ chiếm 10–30 pixel — model bỏ sót nhiều",
         Inches(0.45), CONTENT_Y + Inches(0.08), Inches(12.4), Inches(0.42),
         size=Pt(17), bold=True, color=RED)

# left: steps
add_rect(sl, Inches(0.28), CONTENT_Y + Inches(0.68), Inches(5.6), Inches(3.35), fill=LGRAY)
add_text(sl, "Nguyên lý — Akyon et al., IEEE ICIP 2022",
         Inches(0.38), CONTENT_Y + Inches(0.75), Inches(5.4), Inches(0.42),
         size=Pt(19), bold=True, color=NAVY)
sahi_steps = [
    "① Chia ảnh gốc (640×640) thành các lát nhỏ chồng lấp",
    "② Chạy YOLO inference độc lập trên từng lát ảnh",
    "③ Tổng hợp kết quả từ tất cả lát + ảnh gốc (full scale)",
    "④ NMM — Non-Maximum Merging loại bỏ detection trùng",
]
for i, s in enumerate(sahi_steps):
    y = CONTENT_Y + Inches(1.22 + i * 0.72)
    add_rect(sl, Inches(0.38), y, Inches(5.4), Inches(0.64), fill=WHITE)
    add_rect(sl, Inches(0.38), y, Inches(0.06), Inches(0.64), fill=BLUE)
    add_text(sl, s, Inches(0.52), y + Inches(0.08), Inches(5.1), Inches(0.52),
             size=Pt(17), color=DGRAY, wrap=True)

# SAHI result: landscape cam2 (1.78 ratio) — w=5.6" h=5.6/1.78=3.15" too tall
# Use h=1.55" → w=1.55*1.78=2.76", center in 5.6" wide box
img_border(sl, Inches(0.28), CONTENT_Y + Inches(4.15), Inches(5.6), Inches(1.55), CAM_ANNOTATED,
           "Kết quả sau SAHI — phát hiện Pedestrian & Motorbike rõ hơn")

# right: results
add_rect(sl, Inches(6.05), CONTENT_Y + Inches(0.68), Inches(7.0), Inches(5.1), fill=LGREEN)
add_text(sl, "Kết Quả Thực Nghiệm",
         Inches(6.18), CONTENT_Y + Inches(0.76), Inches(6.8), Inches(0.42),
         size=Pt(17), bold=True, color=GREEN)

add_rect(sl, Inches(6.18), CONTENT_Y + Inches(1.25), Inches(6.7), Inches(1.55), fill=WHITE)
add_text(sl, "Pedestrian Recall",
         Inches(6.3), CONTENT_Y + Inches(1.35), Inches(6.48), Inches(0.38),
         size=Pt(17), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(sl, "0.42  →  0.75",
         Inches(6.3), CONTENT_Y + Inches(1.75), Inches(6.48), Inches(0.65),
         size=Pt(38), bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(sl, "+78.6%  cải thiện với SAHI",
         Inches(6.3), CONTENT_Y + Inches(2.38), Inches(6.48), Inches(0.32),
         size=Pt(17), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_lines(sl, [
    ("Lợi ích trong bài toán fisheye:", True, NAVY, Pt(19)),
    ("• Mỗi lát ảnh có biến dạng cục bộ ít hơn → model dễ học hơn", False, DGRAY, Pt(19)),
    ("• Đối tượng nhỏ ở biên xuất hiện tương đối lớn hơn trong lát", False, DGRAY, Pt(19)),
    ("", False, DGRAY, Pt(11)),
    ("Đánh đổi tốc độ:", True, NAVY, Pt(19)),
    ("• Inference: 1.8–2.5 giây/ảnh (6–8× chậm hơn standard)", False, DGRAY, Pt(19)),
    ("• Phù hợp: phân tích offline, snapshot định kỳ 2–3s/lần", False, DGRAY, Pt(19)),
], Inches(6.18), CONTENT_Y + Inches(2.82), Inches(6.7), Inches(2.85))

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 9 — Bộ Dữ Liệu
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Bộ Dữ Liệu: FishEye8K + VisDrone2019", 9)

# FishEye8K
add_rect(sl, Inches(0.28), CONTENT_Y, Inches(6.05), Inches(3.15), fill=TBLALT)
add_rect(sl, Inches(0.28), CONTENT_Y, Inches(6.05), Inches(0.48), fill=NAVY)
add_text(sl, "FishEye8K — CVPRW 2023, AI City Challenge",
         Inches(0.38), CONTENT_Y + Inches(0.05), Inches(5.85), Inches(0.38),
         size=Pt(19), bold=True, color=WHITE)
add_text(sl, "Camera giám sát Sở Cảnh sát Hsinchu, Đài Loan · Điều kiện đa dạng (ngày/đêm/mưa/sương)",
         Inches(0.38), CONTENT_Y + Inches(0.55), Inches(5.85), Inches(0.52),
         size=Pt(17), color=DGRAY, italic=True, wrap=True)
add_table(sl,
    ["Split", "Số ảnh", "Số nhãn", "TB/ảnh"],
    [["Train","4.230","112.213","21,2"],["Val","1.058","—","—"],["Test","2.712","—","—"],
     [("Tổng",True,WHITE),("8.000",True,WHITE),("112.213",True,WHITE),"—"]],
    x=Inches(0.32), y=CONTENT_Y + Inches(1.12), w=Inches(5.98), h=Inches(1.95),
    col_widths=[Inches(1.4), Inches(1.5), Inches(1.8), Inches(1.28)], fs=Pt(13))

# VisDrone
add_rect(sl, Inches(6.5), CONTENT_Y, Inches(6.55), Inches(3.15), fill=LORANGE)
add_rect(sl, Inches(6.5), CONTENT_Y, Inches(6.55), Inches(0.48), fill=ORANGE)
add_text(sl, "VisDrone2019 — ICCV 2019, Đại học Thiên Tân",
         Inches(6.62), CONTENT_Y + Inches(0.05), Inches(6.35), Inches(0.38),
         size=Pt(19), bold=True, color=WHITE)
add_text(sl, "Ảnh UAV, độ cao 10–70m, góc nghiêng + overhead · Ánh xạ 10 → 5 lớp",
         Inches(6.62), CONTENT_Y + Inches(0.55), Inches(6.35), Inches(0.52),
         size=Pt(17), color=DGRAY, italic=True, wrap=True)
add_table(sl,
    ["Split", "Số ảnh", "Số nhãn"],
    [["Train","6.471","343.205"],["Val","548","38.759"],["Test-Dev","1.610","75.102"]],
    x=Inches(6.54), y=CONTENT_Y + Inches(1.12), w=Inches(6.48), h=Inches(1.95),
    col_widths=[Inches(2.0), Inches(2.2), Inches(2.28)], fs=Pt(13))

# pipeline
add_rect(sl, Inches(0.28), CONTENT_Y + Inches(3.24), Inches(12.77), Inches(0.45), fill=NAVY)
add_text(sl, "Pipeline chuyển đổi VisDrone → Fisheye  (giảm 26% nhãn sau lọc bbox không hợp lệ)",
         Inches(0.42), CONTENT_Y + Inches(3.28), Inches(12.4), Inches(0.38),
         size=Pt(19), bold=True, color=WHITE)

pipe = ["Đọc ảnh\n+ nhãn","to_fisheye()\nstrength=0.5","32-point bbox\ntransform","Lọc bbox\nhợp lệ","Lưu YOLO\nformat"]
for i, s in enumerate(pipe):
    x = Inches(0.28 + i * 2.55)
    add_rect(sl, x, CONTENT_Y + Inches(3.77), Inches(2.38), Inches(1.07),
             fill=BLUE if i%2==0 else LGRAY)
    add_text(sl, s, x + Inches(0.1), CONTENT_Y + Inches(3.87), Inches(2.18), Inches(0.85),
             size=Pt(24), bold=(i%2==0), color=WHITE if i%2==0 else NAVY, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text(sl, "→", Inches(0.28 + i*2.55 + 2.42), CONTENT_Y + Inches(4.1),
                 Inches(0.22), Inches(0.4), size=Pt(24), bold=True, color=NAVY)

# combined + detection image
add_rect(sl, Inches(0.28), CONTENT_Y + Inches(4.96), Inches(7.8), Inches(0.62), fill=LGREEN)
add_text(sl, "Dataset kết hợp:  Train 11.296 ảnh · 406.355 nhãn (TB 35,97 nhãn/ảnh)  |  Val 1.768 ảnh · ~58.000 nhãn",
         Inches(0.4), CONTENT_Y + Inches(5.05), Inches(7.6), Inches(0.55),
         size=Pt(17), bold=True, color=GREEN)

img_border(sl, Inches(8.25), CONTENT_Y + Inches(4.88), Inches(4.78), Inches(0.75), CAM_ANNOTATED,
           "Mẫu detection — camera thực tế")

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 10 — Cấu Hình Huấn Luyện
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Cấu Hình Huấn Luyện", 10)

add_rect(sl, Inches(0.28), CONTENT_Y, Inches(12.77), Inches(0.52), fill=TBLALT)
add_text(sl, "Môi trường: Kaggle Notebooks  ·  GPU Tesla P100-PCIE-16GB  ·  17.1 GB VRAM  ·  RAM 25 GB  ·  CUDA 12.1  ·  PyTorch 2.2",
         Inches(0.42), CONTENT_Y + Inches(0.07), Inches(12.4), Inches(0.42),
         size=Pt(17), color=NAVY, align=PP_ALIGN.CENTER)

add_table(sl,
    ["Siêu tham số", "Cơ Bản (CB)", "Nâng Cao (NC)"],
    [
        ["model",              "yolo11n.pt",       "yolo11n.pt"],
        ["epochs",             "50",               ("80", True, GREEN)],
        ["batch_size",         "16",               "16"],
        ["img_size",           "640",              ("960", True, GREEN)],
        ["optimizer",          "AdamW",            ("SGD + Cosine LR", True, GREEN)],
        ["lr0",                "0.0005",           "0.01"],
        ["freeze backbone",    "Không",            ("Có — 10 lớp đầu", True, GREEN)],
        ["Dataset",            "FishEye8K",        ("FishEye8K + VisDrone", True, GREEN)],
        ["SAHI inference",     "Không",            ("Có", True, GREEN)],
        ["AMP (FP16)",         "True",             "True"],
        ["mosaic / mixup",     "1.0 / 0.05",       "0.8 / 0.15"],
    ],
    x=Inches(0.28), y=CONTENT_Y + Inches(0.61), w=Inches(8.3), h=Inches(5.05),
    col_widths=[Inches(2.8), Inches(2.5), Inches(3.0)], fs=Pt(13))

add_rect(sl, Inches(8.75), CONTENT_Y + Inches(0.61), Inches(4.3), Inches(5.05), fill=LORANGE)
add_text(sl, "Cân bằng dữ liệu\n(Class Imbalance)",
         Inches(8.88), CONTENT_Y + Inches(0.69), Inches(4.1), Inches(0.62),
         size=Pt(17), bold=True, color=ORANGE)
techs = [
    ("Copy-Paste Aug", "Ưu tiên lớp thiểu số\nBus, Truck khi augment"),
    ("Class Weight", "Tự động điều chỉnh\ntheo tần suất lớp"),
    ("Oversampling", "Ưu tiên ảnh chứa\nnhiều Bus / Truck"),
]
for i, (t, d) in enumerate(techs):
    y = CONTENT_Y + Inches(1.4 + i * 1.38)
    add_rect(sl, Inches(8.88), y, Inches(4.0), Inches(1.22), fill=WHITE)
    add_text(sl, t, Inches(8.98), y + Inches(0.1), Inches(3.8), Inches(0.42),
             size=Pt(19), bold=True, color=NAVY)
    add_text(sl, d, Inches(8.98), y + Inches(0.56), Inches(3.8), Inches(0.62),
             size=Pt(17), color=DGRAY, wrap=True)

orange_bar(sl, "Thời gian huấn luyện:  CB ~3.8 giờ  ·  NC ~6.8 giờ (img_size=960, 80 epoch)")

# ══════════════════════════════════════════════════════════
# SLIDE 11 — Kết Quả Chi Tiết Theo Lớp
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Kết Quả Huấn Luyện — Chi Tiết Theo Lớp", 11)

add_text(sl, "CB = Cơ Bản (FishEye8K only)   ·   NC = Nâng Cao (+VisDrone + Freeze + SAHI)",
         Inches(0.28), CONTENT_Y, Inches(12.77), Inches(0.38),
         size=Pt(24), italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

add_table(sl,
    ["Lớp", "Precision\nCB / NC", "Recall\nCB / NC", "mAP@0.5\nCB / NC", "F1\nCB / NC"],
    [
        ["Car",       "0.710 / 0.920", "0.680 / 0.840", "0.720 / 0.910", "0.690 / 0.878"],
        ["Bus",       "0.580 / 0.840", "0.520 / 0.700", "0.590 / 0.820", "0.550 / 0.764"],
        ["Truck",     "0.600 / 0.850", "0.650 / 0.730", "0.610 / 0.830", "0.620 / 0.785"],
        ["Pedestrian","0.720 / 0.830", "0.420 / 0.750", "0.550 / 0.850", "0.530 / 0.788"],
        ["Motorbike", "0.640 / 0.906", "0.580 / 0.790", "0.626 / 0.900", "0.610 / 0.845"],
        [("ALL (mean)",True,WHITE),("0.650 / 0.869",True,WHITE),("0.570 / 0.762",True,WHITE),
         ("0.619 / 0.862",True,WHITE),("0.600 / 0.812",True,WHITE)],
    ],
    x=Inches(0.28), y=CONTENT_Y + Inches(0.45), w=Inches(12.77), h=Inches(3.65),
    col_widths=[Inches(2.3), Inches(2.62), Inches(2.62), Inches(2.62), Inches(2.61)], fs=Pt(13))

hlights = [
    ("mAP@0.5 Cơ Bản",    "0.619", NAVY),
    ("mAP@0.5 Nâng Cao",  "0.862", GREEN),
    ("Cải thiện mAP",     "+39.3%", ORANGE),
    ("Pedestrian Recall", "+78.6%", BLUE),
]
for i, (lbl, val, col) in enumerate(hlights):
    x = Inches(0.28 + i * 3.18)
    add_rect(sl, x, CONTENT_Y + Inches(4.22), Inches(3.07), Inches(1.55), fill=col)
    add_text(sl, val, x + Inches(0.1), CONTENT_Y + Inches(4.35), Inches(2.87), Inches(0.75),
             size=Pt(38), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, x + Inches(0.1), CONTENT_Y + Inches(5.1), Inches(2.87), Inches(0.42),
             size=Pt(17), color=WHITE, align=PP_ALIGN.CENTER)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 12 — So Sánh 2 Phiên Bản
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "So Sánh 2 Phiên Bản & Phân Tích", 12)

add_table(sl,
    ["Phiên bản", "mAP@0.5", "mAP@0.5:0.95", "Precision", "Recall", "FPS (GPU)"],
    [
        ["Cơ Bản (FishEye8K only)", "0.619", "0.363", "0.650", "0.570", "~41 FPS"],
        [("Nâng Cao (+VisDrone\n+SAHI+Freeze)", True, NAVY),
         ("0.862",True,GREEN),("0.572",True,GREEN),("0.869",True,GREEN),("0.762",True,GREEN),"~12 (SAHI)"],
    ],
    x=Inches(0.28), y=CONTENT_Y, w=Inches(12.77), h=Inches(1.52),
    col_widths=[Inches(2.9), Inches(1.95), Inches(2.15), Inches(1.95), Inches(1.95), Inches(1.87)], fs=Pt(13))

improvements = [
    ("mAP@0.5",     "0.619 → 0.862", "+39.3%", GREEN),
    ("mAP@0.5:0.95","0.363 → 0.572", "+57.6%", BLUE),
    ("Precision",   "0.650 → 0.869", "+33.7%", NAVY),
    ("Recall",      "0.570 → 0.762", "+33.7%", ORANGE),
]
for i, (m, chg, pct, col) in enumerate(improvements):
    x = Inches(0.28 + i * 3.18)
    add_rect(sl, x, CONTENT_Y + Inches(1.6), Inches(3.07), Inches(1.55), fill=LGRAY)
    add_rect(sl, x, CONTENT_Y + Inches(1.6), Inches(3.07), Inches(0.06), fill=col)
    add_text(sl, pct,  x + Inches(0.1), CONTENT_Y + Inches(1.68), Inches(2.87), Inches(0.78),
             size=Pt(38), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(sl, m,    x + Inches(0.1), CONTENT_Y + Inches(2.5), Inches(2.87), Inches(0.42),
             size=Pt(17), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(sl, chg,  x + Inches(0.1), CONTENT_Y + Inches(2.88), Inches(2.87), Inches(0.28),
             size=Pt(19), color=DGRAY, align=PP_ALIGN.CENTER)

add_rect(sl, Inches(0.28), CONTENT_Y + Inches(3.24), Inches(12.77), Inches(0.45), fill=NAVY)
add_text(sl, "Phân Tích Yếu Tố Cải Thiện",
         Inches(0.42), CONTENT_Y + Inches(3.28), Inches(12.0), Inches(0.38),
         size=Pt(17), bold=True, color=WHITE)

analyses = [
    ("Freeze backbone\n(freeze=10)", "Bảo toàn đặc trưng pretrained COCO,\ntránh overfitting dataset nhỏ"),
    ("Dataset VisDrone", "Tăng 2× số mẫu → cải thiện\nBus, Truck, Pedestrian đáng kể"),
    ("SAHI inference", "Pedestrian Recall 0.42 → 0.75\n(+78.6% — cải thiện lớn nhất)"),
    ("Đánh đổi tốc độ", "41 FPS → 12 FPS (SAHI)\nChấp nhận được cho offline"),
]
for i, (ttl, dsc) in enumerate(analyses):
    x = Inches(0.28 + i * 3.18)
    add_rect(sl, x, CONTENT_Y + Inches(3.77), Inches(3.07), Inches(2.6), fill=TBLALT)
    add_text(sl, ttl, x + Inches(0.1), CONTENT_Y + Inches(3.85), Inches(2.87), Inches(0.72),
             size=Pt(19), bold=True, color=NAVY)
    add_text(sl, dsc, x + Inches(0.1), CONTENT_Y + Inches(4.62), Inches(2.87), Inches(1.65),
             size=Pt(24), color=DGRAY, wrap=True)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 13 — Kiến Trúc Hệ Thống
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Kiến Trúc Hệ Thống Ứng Dụng", 13)

layers = [
    ("Frontend SPA", "templates/index.html + static/js/*  ·  router.js · appState.js · api.js", RGBColor(0x42, 0x8E, 0xC4)),
    ("Flask App Factory", "create_app(): Config → DB → Blueprints → Extended Routes → Logging → atexit cleanup", NAVY),
    ("Blueprints (routes/)", "core · detect · history · external_camera · examples  +  routes_extended.py: analytics/alerts/speed/congestion/ALPR/export", BLUE),
    ("Module nghiệp vụ", "video_detect · fisheye · speed_estimator · line_counter · congestion_detector · incident_detector · alpr · model_registry", RGBColor(0x1B, 0x5E, 0x20)),
    ("Dữ liệu (db.py)", "SQLite / PostgreSQL  ·  recent_image_store  ·  Dual DB via DATABASE_URL env var", DGRAY),
]
for i, (lbl, dsc, col) in enumerate(layers):
    y = CONTENT_Y + Inches(i * 0.95)
    add_rect(sl, Inches(0.28), y, Inches(7.9), Inches(0.88), fill=col)
    add_text(sl, lbl, Inches(0.38), y + Inches(0.05), Inches(2.2), Inches(0.42),
             size=Pt(19), bold=True, color=WHITE)
    add_text(sl, dsc, Inches(2.68), y + Inches(0.05), Inches(5.3), Inches(0.82),
             size=Pt(14), color=WHITE, wrap=True)
    if i < 4:
        add_text(sl, "▼", Inches(4.0), y + Inches(0.88), Inches(0.5), Inches(0.12),
                 size=Pt(17), color=MGRAY, align=PP_ALIGN.CENTER)

# live detection image (real system output)
img_border(sl, Inches(8.35), CONTENT_Y, Inches(4.9), Inches(4.82), OVERVIEW_DAY,
           "Hệ thống hoạt động thực tế — live 4-camera detection")

add_text(sl, "Flask 3.x · Python 3.10 · YOLOv11 · OpenCV · EasyOCR · SQLite/PostgreSQL · Docker · GCP",
         Inches(0.28), CONTENT_Y + Inches(4.98), Inches(12.77), Inches(0.4),
         size=Pt(17), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 14 — Xử Lý Video Bất Đồng Bộ
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Kiến Trúc Xử Lý Video Bất Đồng Bộ", 14)

add_rect(sl, Inches(0.28), CONTENT_Y, Inches(12.77), Inches(0.58), fill=LRED)
add_rect(sl, Inches(0.28), CONTENT_Y, Inches(0.07), Inches(0.58), fill=RED)
add_text(sl, "Vấn đề: Video 30s, 1080p → ~95 giây xử lý (GTX 1060) — trả lại ngay sẽ gây 504 Gateway Timeout",
         Inches(0.42), CONTENT_Y + Inches(0.08), Inches(12.4), Inches(0.42),
         size=Pt(19), bold=True, color=RED)

add_text(sl, "Giải pháp: Job Queue bất đồng bộ — ThreadPoolExecutor  (HTTP 202 Accepted pattern)",
         Inches(0.28), CONTENT_Y + Inches(0.66), Inches(12.77), Inches(0.4),
         size=Pt(17), bold=True, color=GREEN)

flow = [
    ("Client POST\n/api/detect", BLUE,  "Upload video"),
    ("Validate +\nLưu file tạm",  NAVY, "Server side"),
    ("Tạo job_id\n(UUID v4)",     NAVY, "DB: pending"),
    ("HTTP 202\nAccepted",  GREEN,       "< 200ms ✓"),
    ("Client poll\nGET /jobs/id", BLUE, "Mỗi 2 giây"),
    ("Stream MP4\nkhi 'done'",   GREEN, "Streaming"),
]
for i, (lbl, col, sub) in enumerate(flow):
    x = Inches(0.28 + i * 2.17)
    add_rect(sl, x, CONTENT_Y + Inches(1.14), Inches(2.0), Inches(1.05), fill=col)
    add_text(sl, lbl, x + Inches(0.07), CONTENT_Y + Inches(1.22), Inches(1.86), Inches(0.75),
             size=Pt(19), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, sub, x + Inches(0.07), CONTENT_Y + Inches(1.97), Inches(1.86), Inches(0.32),
             size=Pt(14), color=MGRAY, align=PP_ALIGN.CENTER)
    if i < 5:
        add_text(sl, "→", Inches(0.28 + i*2.17 + 2.03), CONTENT_Y + Inches(1.48),
                 Inches(0.22), Inches(0.38), size=Pt(24), bold=True, color=NAVY)

add_rect(sl, Inches(0.28), CONTENT_Y + Inches(2.35), Inches(12.77), Inches(2.72), fill=LGRAY)
add_rect(sl, Inches(0.28), CONTENT_Y + Inches(2.35), Inches(0.07), Inches(2.72), fill=ORANGE)
add_text(sl, "Background Worker Thread — VideoJobQueue.submit()",
         Inches(0.45), CONTENT_Y + Inches(2.42), Inches(12.0), Inches(0.45),
         size=Pt(17), bold=True, color=NAVY)

worker = [
    "① Đọc từng frame video (OpenCV)", "② YOLO inference trên frame",
    "③ Speed + Congestion overlay",     "④ Ghi annotation vào frame",
    "⑤ Ghép annotated.mp4 + preview",  "⑥ Cập nhật DB: done / failed",
]
for i, s in enumerate(worker):
    col = i % 3; row = i // 3
    x = Inches(0.5 + col * 4.2); y = CONTENT_Y + Inches(2.95 + row * 0.62)
    add_rect(sl, x, y, Inches(0.32), Inches(0.32), fill=ORANGE)
    add_text(sl, str(i+1), x, y, Inches(0.32), Inches(0.32),
             size=Pt(19), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, s, x + Inches(0.4), y, Inches(3.72), Inches(0.38), size=Pt(17), color=DGRAY)

add_rect(sl, Inches(0.28), CONTENT_Y + Inches(5.17), Inches(12.77), Inches(0.58), fill=TBLALT)
add_text(sl, "max_workers=2  ·  max_queue_size=10  ·  Job states: pending → running → done / failed  ·  Tránh quá tải VRAM",
         Inches(0.42), CONTENT_Y + Inches(5.26), Inches(12.4), Inches(0.42),
         size=Pt(19), color=NAVY, align=PP_ALIGN.CENTER)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 15 — 5 Module ITS
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "5 Module Phân Tích Giao Thông ITS", 15)

modules = [
    ("01","SpeedEstimator",      "IoU tracking liên frame → pixel displacement → km/h · Cảnh báo vi phạm tốc độ", NAVY),
    ("02","CongestionDetector",  "ROI mật độ → 4 mức Level of Service · Visualize overlay màu theo mức độ", BLUE),
    ("03","IncidentDetector",    "3 loại sự cố: đỗ sai · dừng đột ngột · ngược chiều · Alert webhook async", RED),
    ("04","LineCounter",         "4 vạch ảo N/S/E/W · Cross-product xác định hướng · Đếm xe theo giờ", ORANGE),
    ("05","ALPR",                "YOLO crop → EasyOCR → Regex chuẩn biển VN → Lưu DB + tra cứu lịch sử", GREEN),
]
for i, (num, ttl, dsc, col) in enumerate(modules):
    y = CONTENT_Y + Inches(i * 1.0)
    add_rect(sl, Inches(0.28), y, Inches(7.62), Inches(0.93), fill=LGRAY)
    add_rect(sl, Inches(0.28), y, Inches(0.9), Inches(0.93), fill=col)
    add_text(sl, num, Inches(0.28), y, Inches(0.9), Inches(0.93),
             size=Pt(36), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, ttl, Inches(1.28), y + Inches(0.07), Inches(3.0), Inches(0.42),
             size=Pt(17), bold=True, color=col)
    add_text(sl, dsc, Inches(1.28), y + Inches(0.52), Inches(6.5), Inches(0.48),
             size=Pt(14), color=DGRAY, wrap=True)

img_border(sl, Inches(8.08), CONTENT_Y, Inches(5.18), Inches(5.02), OVERVIEW_DAY,
           "Live detection — 4 camera, real-time")

add_text(sl, "Hỗ trợ chung: alert_manager (webhook) · analytics (heatmap + hourly) · cloud_storage (GCS)",
         Inches(0.28), CONTENT_Y + Inches(5.12), Inches(12.77), Inches(0.38),
         size=Pt(24), italic=True, color=BLUE, align=PP_ALIGN.CENTER)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 16 — SpeedEstimator
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Module SpeedEstimator — IoU Tracking", 16)

add_rect(sl, Inches(0.28), CONTENT_Y, Inches(6.1), Inches(5.85), fill=LGRAY)
add_text(sl, "Thuật Toán IoU Tracking",
         Inches(0.38), CONTENT_Y + Inches(0.08), Inches(5.9), Inches(0.42),
         size=Pt(17), bold=True, color=NAVY)
add_text(sl, "(Đơn giản hơn SORT/DeepSORT — không cần model re-ID)",
         Inches(0.38), CONTENT_Y + Inches(0.53), Inches(5.9), Inches(0.32),
         size=Pt(17), color=DGRAY, italic=True)
add_rect(sl, Inches(0.38), CONTENT_Y + Inches(0.93), Inches(5.9), Inches(1.78), fill=NAVY)
add_text(sl, 'Frame t:   D = {(bbox, class, conf)}\nFrame t+1: Ma trận IoU giữa tất cả cặp\n           Hungarian matching → max Σ IoU\n           threshold=0.3  ·  max_age=5 frames\n           Unmatched → new track; Lost → xóa',
         Inches(0.48), CONTENT_Y + Inches(1.02), Inches(5.7), Inches(1.6),
         size=Pt(11), color=RGBColor(0x90, 0xFF, 0x90), font='Consolas')

add_text(sl, "Công thức tốc độ:",
         Inches(0.38), CONTENT_Y + Inches(2.82), Inches(5.9), Inches(0.38),
         size=Pt(19), bold=True, color=NAVY)
add_rect(sl, Inches(0.38), CONTENT_Y + Inches(3.28), Inches(5.9), Inches(1.35), fill=WHITE)
add_text(sl, "v (m/s) = √(Δcx² + Δcy²) × fps / pixels_per_meter\nv (km/h) = v × 3.6\npixels_per_meter = 8.0 px/m  (camera H=5m, góc 45°)",
         Inches(0.48), CONTENT_Y + Inches(3.38), Inches(5.7), Inches(1.15),
         size=Pt(12), color=NAVY, font='Consolas')
add_text(sl, "fisheye_correction=True: hệ số bù trừ biến dạng hướng kính vùng biên",
         Inches(0.38), CONTENT_Y + Inches(4.72), Inches(5.9), Inches(0.5),
         size=Pt(17), color=DGRAY, italic=True, wrap=True)

speed_levels = [
    ("< 40 km/h",  GREEN, "Xanh lá — Bình thường"),
    ("40–70 km/h", ORANGE, "Cam — Nhanh"),
    ("> 70 km/h",  RED,    "Đỏ — ⚠ Vượt tốc"),
]
for i, (sp, col, lbl) in enumerate(speed_levels):
    y = CONTENT_Y + Inches(i * 1.52)
    add_rect(sl, Inches(6.55), y, Inches(6.5), Inches(1.42), fill=col)
    add_text(sl, sp,  Inches(6.7), y + Inches(0.12), Inches(4.0), Inches(0.62),
             size=Pt(36), bold=True, color=WHITE)
    add_text(sl, lbl, Inches(6.7), y + Inches(0.78), Inches(6.1), Inches(0.48),
             size=Pt(17), color=WHITE)

add_text(sl, "Cảnh báo vi phạm: Tốc độ > ngưỡng trong ≥ 3 frame liên tiếp → lưu bản ghi vi phạm",
         Inches(6.55), CONTENT_Y + Inches(4.62), Inches(6.5), Inches(0.72),
         size=Pt(19), bold=True, color=NAVY, wrap=True)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 17 — CongestionDetector
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Module CongestionDetector — Level of Service", 17)

add_rect(sl, Inches(0.28), CONTENT_Y, Inches(12.77), Inches(0.62), fill=TBLALT)
add_text(sl, "Phương pháp: Phân tích mật độ trong ROI (tọa độ chuẩn hóa ∈ [0,1]²)   "
             "·   Trọng số: Motorbike 0.5 · Car 1.0 · Truck 2.0 · Bus 2.5 · Pedestrian 0.3",
         Inches(0.42), CONTENT_Y + Inches(0.1), Inches(12.4), Inches(0.48),
         size=Pt(17), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

los = [
    ("FREE\nThông thoáng",      "density < 0.3",       "Luồng giao thông\nbình thường",          GREEN),
    ("MODERATE\nVừa phải",      "0.3 ≤ density < 0.6", "Lưu lượng TB\nkhông cần can thiệp",     RGBColor(0xF5, 0x7F, 0x17)),
    ("HEAVY\nNặng",             "0.6 ≤ density < 0.9", "Tắc nghẽn cục bộ\ncần chú ý theo dõi",  ORANGE),
    ("SEVERE\nNghiêm trọng",    "density ≥ 0.9",        "Tắc nghẽn nghiêm trọng\nCần can thiệp ngay", RED),
]
for i, (lvl, cond, meaning, col) in enumerate(los):
    x = Inches(0.28 + i * 3.18)
    add_rect(sl, x, CONTENT_Y + Inches(0.72), Inches(3.07), Inches(3.78), fill=col)
    add_text(sl, lvl, x + Inches(0.1), CONTENT_Y + Inches(0.82), Inches(2.87), Inches(1.05),
             size=Pt(24), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, x + Inches(0.12), CONTENT_Y + Inches(1.92), Inches(2.83), Inches(0.55), fill=WHITE)
    add_text(sl, cond, x + Inches(0.12), CONTENT_Y + Inches(2.0), Inches(2.83), Inches(0.42),
             size=Pt(13), bold=True, color=col, align=PP_ALIGN.CENTER, font='Consolas')
    add_text(sl, meaning, x + Inches(0.1), CONTENT_Y + Inches(2.58), Inches(2.87), Inches(1.72),
             size=Pt(24), color=WHITE, align=PP_ALIGN.CENTER, wrap=True)

add_rect(sl, Inches(0.28), CONTENT_Y + Inches(4.6), Inches(12.77), Inches(0.55), fill=LGRAY)
add_text(sl, "Visualization: Hình chữ nhật bán trong suốt theo màu  ·  Text: tên ROI, số xe/capacity, %  ·  Dashboard tổng hợp",
         Inches(0.42), CONTENT_Y + Inches(4.67), Inches(12.4), Inches(0.42),
         size=Pt(17), color=NAVY)

add_rect(sl, Inches(0.28), CONTENT_Y + Inches(5.22), Inches(12.77), Inches(0.45), fill=RED)
add_text(sl, "Alert: SEVERE → Gửi webhook (không block pipeline chính)  ·  alert_manager async call",
         Inches(0.42), CONTENT_Y + Inches(5.28), Inches(12.4), Inches(0.38),
         size=Pt(17), bold=True, color=WHITE)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 18 — Analytics & ALPR
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Module Analytics, Heatmap & ALPR", 18)

add_rect(sl, Inches(0.28), CONTENT_Y, Inches(6.1), Inches(5.85), fill=LGRAY)
add_text(sl, "Heatmap Mật Độ (analytics.py)",
         Inches(0.38), CONTENT_Y + Inches(0.08), Inches(5.9), Inches(0.42),
         size=Pt(17), bold=True, color=NAVY)
add_rect(sl, Inches(0.38), CONTENT_Y + Inches(0.57), Inches(5.9), Inches(1.7), fill=NAVY)
add_text(sl, "# Mỗi detection tại (cx, cy):\nheatmap_acc[cy, cx] += 1\n\n# Render định kỳ:\nGaussian blur(sigma=15) → Normalize [0,255]\n→ cv2.COLORMAP_INFERNO\n→ Blend ảnh gốc (alpha=0.4)",
         Inches(0.48), CONTENT_Y + Inches(0.65), Inches(5.7), Inches(1.55),
         size=Pt(11), color=RGBColor(0x90, 0xFF, 0x90), font='Consolas')

add_text(sl, "Dashboard Analytics:",
         Inches(0.38), CONTENT_Y + Inches(2.4), Inches(5.9), Inches(0.38),
         size=Pt(19), bold=True, color=NAVY)
dashboard = [
    "• Biểu đồ lưu lượng theo giờ (bar chart)",
    "• Phân bố loại phương tiện (pie chart)",
    "• Xác định giờ cao điểm tự động",
    "• TOC Traffic Operations Center — 8 trang SPA",
    "• Export CSV / JSON",
]
add_lines(sl, [(s, False, DGRAY, Pt(19)) for s in dashboard],
          Inches(0.38), CONTENT_Y + Inches(2.85), Inches(5.9), Inches(2.85))

# right: ALPR
add_rect(sl, Inches(6.55), CONTENT_Y, Inches(6.5), Inches(5.85), fill=LORANGE)
add_text(sl, "Module ALPR — Nhận Dạng Biển Số VN",
         Inches(6.68), CONTENT_Y + Inches(0.08), Inches(6.3), Inches(0.5),
         size=Pt(17), bold=True, color=ORANGE)
alpr_steps = [
    "① YOLO bbox → Crop vùng xe (ưu tiên xe lớn nhất gần camera)",
    '② EasyOCR → raw text',
    '③ Regex chuẩn hóa biển VN:\n   "51F-123.45"  hoặc  "29X1-2345"',
    "④ Validate: confidence ≥ 0.30 · độ dài 7–10 ký tự",
    "⑤ Lưu bảng license_plates · tra cứu lịch sử phương tiện",
]
for i, s in enumerate(alpr_steps):
    y = CONTENT_Y + Inches(0.67 + i * 0.92)
    add_rect(sl, Inches(6.68), y, Inches(0.42), Inches(0.42), fill=ORANGE)
    add_text(sl, str(i+1), Inches(6.68), y, Inches(0.42), Inches(0.42),
             size=Pt(19), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, s, Inches(7.2), y, Inches(5.7), Inches(0.85),
             size=Pt(19), color=DGRAY, wrap=True)

img_border(sl, Inches(6.68), CONTENT_Y + Inches(5.28), Inches(6.1), Inches(0.5), CAM_ANNOTATED,
           "ALPR — phát hiện & nhận dạng biển số thực tế")

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 19 — Giao Diện Web
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Giao Diện Web — 6 Phân Hệ", 19)

add_text(sl, "Kiến trúc Frontend: Vanilla JavaScript SPA — HTML5 · Bootstrap 5 · ES6 Modules",
         Inches(0.28), CONTENT_Y, Inches(12.77), Inches(0.38),
         size=Pt(19), italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

add_table(sl,
    ["Phân hệ", "Chức năng chi tiết"],
    [
        [("Dashboard",           True, NAVY), "KPIs real-time: tổng lượt chạy, số đối tượng, biểu đồ phân phối, system logs"],
        [("Workspace (Inference)",True, NAVY), "Kéo-thả ảnh/video · Đồng bộ (ảnh <500ms) · Bất đồng bộ job queue (video)"],
        [("Live Streams",         True, NAVY), "MJPEG stream · Song song: camera gốc + detection overlay + tốc độ + ùn tắc"],
        [("Run History",          True, NAVY), "Lịch sử phiên nhận diện · Tải ảnh/video annotated + metadata JSON"],
        [("System Logs",          True, NAVY), "Terminal view · Server-Sent Events (SSE) · Tìm kiếm, lọc theo log level"],
        [("TOC / ALPR",           True, NAVY), "Trung tâm điều hành · Nhận dạng & tra cứu biển số · Heatmap · Export"],
    ],
    x=Inches(0.28), y=CONTENT_Y + Inches(0.45), w=Inches(8.25), h=Inches(4.0),
    col_widths=[Inches(2.8), Inches(5.45)], fs=Pt(13))

img_border(sl, Inches(8.62), CONTENT_Y + Inches(0.45), Inches(4.73), Inches(3.95), OVERVIEW_DAY,
           "Live stream — 4 camera đồng thời")

# UX principles
ux = [
    ("Responsive", "Desktop 1920×1080\n+ Tablet 1024×768"),
    ("Progressive\nDisclosure", "Chi tiết kỹ thuật\nẩn mặc định"),
    ("Error UX", '"File quá lớn.\nVui lòng chọn < 500MB"'),
    ("Feedback", "Loading spinner\nToast notification"),
]
for i, (ttl, dsc) in enumerate(ux):
    x = Inches(0.28 + i * 3.18)
    add_rect(sl, x, CONTENT_Y + Inches(4.55), Inches(3.07), Inches(1.25), fill=TBLALT)
    add_rect(sl, x, CONTENT_Y + Inches(4.55), Inches(3.07), Inches(0.06), fill=BLUE)
    add_text(sl, ttl, x + Inches(0.1), CONTENT_Y + Inches(4.63), Inches(2.87), Inches(0.42),
             size=Pt(19), bold=True, color=NAVY)
    add_text(sl, dsc, x + Inches(0.1), CONTENT_Y + Inches(5.1), Inches(2.87), Inches(0.65),
             size=Pt(17), color=DGRAY, wrap=True)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 20 — Kiểm Thử
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Kết Quả Kiểm Thử Chức Năng", 20)

add_table(sl,
    ["ID", "Kịch bản kiểm thử", "Kết quả", "Thời gian"],
    [
        ["TC-01", "Upload ảnh JPEG 1920×1080, phát hiện đối tượng",          "✅ PASS", "480ms"],
        ["TC-02", "Upload ảnh > 50MB (vượt giới hạn kích thước)",             "✅ PASS (Error 413)", "50ms"],
        ["TC-03", "Upload video 30s 1080p, xử lý bất đồng bộ",               "✅ PASS (Job async)", "Tạo: 200ms"],
        ["TC-04", "Polling job status khi đang xử lý (running)",               "✅ PASS (status + %)", "45ms"],
        ["TC-05", "Download video kết quả đã annotate",                        "✅ PASS (Stream MP4)", "Streaming"],
        ["TC-06", "SAHI inference trên ảnh đông người đi bộ",                  "✅ PASS", "2.1s"],
        ["TC-07", "API /api/health check",                                     "✅ PASS", "12ms"],
        ["TC-08", "Gửi webhook khi mật độ SEVERE",                             "✅ PASS", "~100ms"],
        ["TC-09", "Concurrent 3 video jobs đồng thời",                        "✅ PASS (Queue thứ tự)", "—"],
    ],
    x=Inches(0.28), y=CONTENT_Y, w=Inches(12.77), h=Inches(5.3),
    col_widths=[Inches(1.28), Inches(6.4), Inches(3.0), Inches(2.09)], fs=Pt(13))

add_rect(sl, Inches(0.28), CONTENT_Y + Inches(5.4), Inches(12.77), Inches(0.68), fill=LGREEN)
add_rect(sl, Inches(0.28), CONTENT_Y + Inches(5.4), Inches(0.07), Inches(0.68), fill=GREEN)
add_text(sl, "Tổng kết:  9/9 Functional Tests PASS   ·   Unit Tests (pytest): 31/31 PASS",
         Inches(0.45), CONTENT_Y + Inches(5.55), Inches(12.4), Inches(0.45),
         size=Pt(20), bold=True, color=GREEN, align=PP_ALIGN.CENTER)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 21 — Hiệu Năng Tổng Thể
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Đánh Giá Hiệu Năng Tổng Thể", 21)

add_text(sl, "Môi trường deployment thực tế: NVIDIA GTX 1060 6GB GDDR5",
         Inches(0.28), CONTENT_Y, Inches(12.77), Inches(0.38),
         size=Pt(19), bold=True, italic=True, color=DGRAY, align=PP_ALIGN.CENTER)

add_table(sl,
    ["Tác vụ", "Thời gian thực đo", "Yêu cầu / Ghi chú"],
    [
        ["Phát hiện ảnh đơn (1080p)",       "380–520ms  (TB: 450ms)", "✅ Đáp ứng yêu cầu < 500ms"],
        ["Xử lý video 1080p, 25fps, 30s",   "~95 giây (3.2× realtime)", "~41 FPS effective (P100 Kaggle)"],
        ["SAHI inference",                   "1.8–2.5s/ảnh",           "Recall người đi bộ 32% → 58%"],
        ["RAM sử dụng",                     "~1.8 GB idle / ~3.2 GB",  "Môi trường GTX 1060"],
        ["VRAM sử dụng",                    "~2.1 GB (YOLOv11 FP16)",  "—"],
        ["API list-jobs",                    "< 100ms",                 "—"],
    ],
    x=Inches(0.28), y=CONTENT_Y + Inches(0.47), w=Inches(12.77), h=Inches(3.05),
    col_widths=[Inches(3.5), Inches(3.8), Inches(5.47)], fs=Pt(13))

add_rect(sl, Inches(0.28), CONTENT_Y + Inches(3.63), Inches(12.77), Inches(0.45), fill=NAVY)
add_text(sl, "Yêu Cầu Phi Chức Năng Đạt Được",
         Inches(0.42), CONTENT_Y + Inches(3.67), Inches(12.0), Inches(0.38),
         size=Pt(17), bold=True, color=WHITE)

nfrs = [
    "✅ Xử lý ảnh ≤ 500ms (GPU)",
    "✅ API submit video < 200ms",
    "✅ Uptime ≥ 99% giờ cao điểm",
    "✅ Dual DB (SQLite/PostgreSQL)",
    "✅ Chrome 90+, Firefox 88+, Edge 90+",
]
for i, nfr in enumerate(nfrs):
    x = Inches(0.28 + i * 2.52)
    add_rect(sl, x, CONTENT_Y + Inches(4.16), Inches(2.42), Inches(1.6), fill=LGREEN)
    add_text(sl, nfr, x + Inches(0.1), CONTENT_Y + Inches(4.25), Inches(2.25), Inches(1.35),
             size=Pt(17), bold=True, color=GREEN, wrap=True)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 22 — Kết Quả Đạt Được
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Kết Quả Đạt Được", 22)

add_rect(sl, Inches(0.28), CONTENT_Y, Inches(6.05), Inches(5.85), fill=LGRAY)
add_rect(sl, Inches(0.28), CONTENT_Y, Inches(6.05), Inches(0.52), fill=NAVY)
add_text(sl, "Phần 1 — Nghiên Cứu Mô Hình",
         Inches(0.38), CONTENT_Y + Inches(0.07), Inches(5.85), Inches(0.4),
         size=Pt(17), bold=True, color=WHITE)
model_results = [
    "✅ Pipeline VisDrone → fisheye với 32-point bbox transform",
    "✅ Dataset kết hợp: 11.296 ảnh train · 406.355 nhãn",
    "✅ YOLOv11-N Cơ bản:   mAP@0.5 = 0.619",
    "✅ YOLOv11-N Nâng cao: mAP@0.5 = 0.862  (+39.3%)",
    "✅ SAHI: Pedestrian Recall 0.42 → 0.75  (+78.6%)",
    "✅ Chỉ 2.6M tham số — phù hợp edge deployment",
]
add_lines(sl, [(r, False, GREEN if "✅" in r else DGRAY, Pt(19)) for r in model_results],
          Inches(0.38), CONTENT_Y + Inches(0.62), Inches(5.85), Inches(4.85), spacing=1.4)

add_rect(sl, Inches(6.5), CONTENT_Y, Inches(6.55), Inches(2.98), fill=TBLALT)
add_rect(sl, Inches(6.5), CONTENT_Y, Inches(6.55), Inches(0.52), fill=BLUE)
add_text(sl, "Phần 2 — Ứng Dụng Giám Sát",
         Inches(6.62), CONTENT_Y + Inches(0.07), Inches(6.35), Inches(0.4),
         size=Pt(17), bold=True, color=WHITE)
app_results = [
    "✅ Flask Application Factory + Blueprint modular",
    "✅ Xử lý video bất đồng bộ — không timeout",
    "✅ 5 module ITS: Speed · Congestion · Incident · LineCounter · ALPR",
    "✅ Dual DB, Cloud Storage (GCS), CI/CD → GCP (NVIDIA L4)",
    "✅ 9/9 functional tests · 31/31 unit tests PASS",
]
add_lines(sl, [(r, False, GREEN if "✅" in r else DGRAY, Pt(19)) for r in app_results],
          Inches(6.62), CONTENT_Y + Inches(0.62), Inches(6.35), Inches(2.2), spacing=1.35)

img_border(sl, Inches(6.5), CONTENT_Y + Inches(3.05), Inches(6.55), Inches(2.78), OVERVIEW_NIGHT,
           "Live detection ban đêm — 4 camera, xe máy Việt Nam")

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 23 — Hạn Chế & Hướng Phát Triển
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Hạn Chế & Hướng Phát Triển", 23)

add_table(sl,
    ["Hạn chế", "Chi tiết"],
    [
        ["Recall Pedestrian thấp (CB)", "Recall 0.42 ở phiên bản Cơ Bản — SAHI cải thiện lên 0.75 nhưng cần thêm dữ liệu đa dạng"],
        ["Calibrate tốc độ thủ công",  "pixels_per_meter=8.0 đặt tay — cần quy trình camera calibration tự động từ điểm mốc thực"],
        ["ROI tắc nghẽn cố định",      "User phải cấu hình ROI tay — kém linh hoạt khi camera dịch chuyển hoặc thay đổi góc"],
        ["Dữ liệu chưa đặc thù VN",   "FishEye8K (Đài Loan) + VisDrone (Trung Quốc) — xe máy và phong cách lái VN khác biệt đáng kể"],
        ["Chưa có edge deployment",    "Chạy server tập trung — chưa tối ưu cho Jetson Nano, Raspberry Pi (< 25 FPS target)"],
    ],
    x=Inches(0.28), y=CONTENT_Y, w=Inches(12.77), h=Inches(2.58),
    col_widths=[Inches(3.0), Inches(9.77)], fs=Pt(13))

add_rect(sl, Inches(0.28), CONTENT_Y + Inches(2.68), Inches(12.77), Inches(0.45), fill=ORANGE)
add_text(sl, "Hướng Phát Triển",
         Inches(0.42), CONTENT_Y + Inches(2.72), Inches(12.0), Inches(0.38),
         size=Pt(17), bold=True, color=WHITE)

future = [
    ("Thu thập\ndữ liệu VN", "2–3 nút giao Hà Nội\n~5.000 ảnh, 3 tháng\nMotorbike + Pedestrian", NAVY),
    ("Kiến trúc\nmới", "YOLOv12 · RT-DETR\nSelf-attention toàn cục\nXử lý distortion tốt hơn", BLUE),
    ("Edge\ncomputing", "TensorRT/ONNX Runtime\nJetson Orin Nano\nMục tiêu ≥ 25 FPS", RED),
    ("Dashboard\ncải tiến", "React.js + WebSocket\nOpenLayers map real-time\nBáo cáo PDF/Excel tự động", GREEN),
]
for i, (ttl, dsc, col) in enumerate(future):
    x = Inches(0.28 + i * 3.18)
    add_rect(sl, x, CONTENT_Y + Inches(3.2), Inches(3.07), Inches(2.97), fill=LGRAY)
    add_rect(sl, x, CONTENT_Y + Inches(3.2), Inches(3.07), Inches(0.72), fill=col)
    add_text(sl, ttl, x + Inches(0.1), CONTENT_Y + Inches(3.28), Inches(2.87), Inches(0.65),
             size=Pt(17), bold=True, color=WHITE)
    add_text(sl, dsc, x + Inches(0.1), CONTENT_Y + Inches(3.97), Inches(2.87), Inches(2.1),
             size=Pt(17), color=DGRAY, wrap=True)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SLIDE 24 — Kết Luận
# ══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
add_bg(sl)
slide_title(sl, "Kết Luận", 24)

add_rect(sl, Inches(0.28), CONTENT_Y, Inches(6.05), Inches(5.58), fill=WHITE,
         line=NAVY, line_w=Pt(1.5))
add_rect(sl, Inches(0.28), CONTENT_Y, Inches(6.05), Inches(0.55), fill=BLUE)
add_text(sl, "Về Nghiên Cứu Mô Hình",
         Inches(0.38), CONTENT_Y + Inches(0.08), Inches(5.85), Inches(0.42),
         size=Pt(17), bold=True, color=WHITE)
conc1 = [
    "• Pipeline 32-point bbox transform VisDrone → fisheye",
    "• Fine-tune YOLOv11-N:  mAP@0.5 = 0.862  (+39.3%)",
    "• SAHI: Recall người đi bộ  0.42 → 0.75  (+78.6%)",
    "• 2.6M params, 6.5 GFLOPs — phù hợp edge deployment",
]
add_lines(sl, [(r, False, NAVY, Pt(17)) for r in conc1],
          Inches(0.38), CONTENT_Y + Inches(0.68), Inches(5.85), Inches(4.55), spacing=1.5)

add_rect(sl, Inches(6.5), CONTENT_Y, Inches(6.55), Inches(5.58), fill=WHITE,
         line=ORANGE, line_w=Pt(1.5))
add_rect(sl, Inches(6.5), CONTENT_Y, Inches(6.55), Inches(0.55), fill=ORANGE)
add_text(sl, "Về Ứng Dụng & Thực Tiễn",
         Inches(6.62), CONTENT_Y + Inches(0.08), Inches(6.35), Inches(0.42),
         size=Pt(17), bold=True, color=WHITE)
conc2 = [
    "• 5 module ITS: Speed · Congestion · Incident · Line · ALPR",
    "• 9/9 functional tests · 31/31 unit tests PASS",
    "• Xử lý ảnh < 500ms — sẵn sàng production deployment",
    "• 1 fisheye = 4 camera thường tại ngã tư đô thị",
]
add_lines(sl, [(r, False, NAVY, Pt(17)) for r in conc2],
          Inches(6.62), CONTENT_Y + Inches(0.68), Inches(6.35), Inches(2.58), spacing=1.5)

img_border(sl, Inches(6.62), CONTENT_Y + Inches(3.35), Inches(6.28), Inches(2.15), CAM_ANNOTATED,
           "Hệ thống nhận diện thực tế — camera Việt Nam")

add_rect(sl, Inches(0.28), CONTENT_Y + Inches(5.65), Inches(12.77), Inches(0.62), fill=NAVY)
add_text(sl, "Đề tài góp phần xây dựng nền tảng ITS Việt Nam ứng dụng AI bản địa",
         Inches(0.42), CONTENT_Y + Inches(5.78), Inches(12.4), Inches(0.42),
         size=Pt(19), bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

orange_bar(sl)

# ══════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════
out = r"C:\Using\DATN\BaoVe_NguyenQuocNam_221220938.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
