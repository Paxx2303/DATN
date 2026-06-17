from pptx import Presentation

prs = Presentation('BaoVe_NguyenQuocNam_221220938.pptx')
print(f'Slides: {len(prs.slides)}')
print(f'Width: {prs.slide_width.inches:.2f}"  Height: {prs.slide_height.inches:.2f}"')
print()

for idx in [0, 1, 6, 8, 16, 22, 23]:
    sl = prs.slides[idx]
    texts = []
    for shape in sl.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t and len(t) > 3:
                    texts.append(t[:55])
    print(f'Slide {idx+1} ({len(sl.shapes)} shapes): {texts[:4]}')
